from django.shortcuts import render

# 显示的时间是带时区的（自动）
from django.utils import timezone
from django.forms.models import model_to_dict
from django.conf import settings
from django.contrib.auth.decorators import login_required
from django.http import HttpResponse, HttpResponseRedirect, JsonResponse, response

# 使用到的数据表要先导入
from database.models import Book
from database.models import BookHistory
from database.models import BookPage
from database.models import BookPageHistory
from database.models import TrainingRecord, QuestionOption, ChoiceQuestionHistory
from database.models import (
    User,
    Technology,
    Plant,
    QuestionBank,
    TrainingBank,
    ChoiceQuestion,
    Examinee,
    ExamPaper,
    ExamedPaper,
    ExamedChoiceQuestion,
    ExamedQuestionOption,
    Option,
    Examinee,
)
from database.models import Option
from django.db.models import Count

from django.views.decorators.csrf import csrf_exempt
from datetime import datetime
from django.contrib.auth.hashers import check_password
import json
from django.contrib import auth
from database import models
from django.core import serializers
from django.db import connection
from django.db.models import F
from django.db.models import Q

# 读.pptx的信息
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from collections import Counter

import re
import logging
import psycopg2
import zipfile
import os
import glob
import shutil

import numpy as np
import pandas as pd
import random

from django.urls import reverse

import time

# Create your views here.
# def test(self, request):
#         return render(request, 'website/templates/test.html')

TRAINING_BANK = "training_bank"
QUESTION_BANK = "question_bank"

logger = logging.getLogger()


def SuperAdminIndex(request):
    return render(request, "superAdminIndex.html")


# 超级管理员获取全部users
@csrf_exempt
def get_super_all_user(request):
    try:
        UserList = serializers.serialize(
            "json", User.objects.all().order_by('Site', 'PlantName', 'StaffID')
        )
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": UserList})


# 超级管理员查找所有site
@csrf_exempt
def get_super_all_site_plant(request):
    try:
        siteList = Plant.objects.all().values().order_by('Site', 'ShortName')
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({'siteList': list(siteList)})


# 超级管理员添加厂别
@csrf_exempt
def add_plant(request):
    new_site = request.POST.get('new_site')
    exist_site = request.POST.get('exist_site')
    plant_ipt = request.POST.get('plant_ipt')
    lang_sel = request.POST.get('lang_sel')
    try:
        if new_site:
            newsite = Plant(Site=new_site, ShortName=plant_ipt, Language=lang_sel)
            newsite.save()
        if exist_site:
            newExistSiteObj = Plant(
                Site=exist_site, ShortName=plant_ipt, Language=lang_sel
            )
            newExistSiteObj.save()
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"res": '添加厂别成功！'})


# 超级管理员修改Plant
@csrf_exempt
def modify_plant(request):
    palnt_id = request.POST.get('palnt_id')
    modifyLanguage = request.POST.get('modifyLanguage')
    PlantObj = Plant.objects.get(PlantID=palnt_id)
    PlantObj.Language = modifyLanguage
    PlantObj.save()
    return JsonResponse({"data": "true"})


# 超级管理员删除厂别
@csrf_exempt
def del_plant(request):
    palnt_id = request.POST.get('palnt_id')
    plant_del = Plant.objects.get(PlantID=palnt_id)
    plant_del.delete()
    return JsonResponse({"data": "true"})


def display(req):
    return render(req, "test.html")


def index(request):
    return render(request, "base.html")


@csrf_exempt
def get_user_authority(request):
    user_id = request.POST.get('Staff_ID', 0)

    is_training_material_maker = "false"
    is_trainer = "false"
    is_exam_question_maker = "false"
    is_exam_paper_publisher = "false"

    training_bank_objs = TrainingBank.objects.filter(
        TrainingMaterialMaker__contains=user_id
    )
    if training_bank_objs.exists():
        is_training_material_maker = "true"

    training_bank_objs = TrainingBank.objects.filter(Trainer__contains=user_id)
    if training_bank_objs.exists():
        is_trainer = "true"

    question_bank_objs = QuestionBank.objects.filter(
        ExamQuestionMaker__contains=user_id
    )
    if question_bank_objs.exists():
        is_exam_question_maker = "true"

    question_bank_objs = QuestionBank.objects.filter(
        ExamPaperPublisher__contains=user_id
    )
    if question_bank_objs.exists():
        is_exam_paper_publisher = "true"

    context = {
        'is_training_material_maker': is_training_material_maker,
        'is_trainer': is_trainer,
        'is_exam_question_maker': is_exam_question_maker,
        'is_exam_paper_publisher': is_exam_paper_publisher,
    }

    return JsonResponse(context)


@csrf_exempt
def index_todo_num(request):
    user_id = request.POST.get('Staff_ID', '')

    to_exam_paper_num = Examinee.objects.filter(StaffID=user_id, IsPassed=False).count()
    exam_paper_obj_list = ExamPaper.objects.filter(CreatorID=user_id)
    examinee_obj_list = (
        Examinee.objects.filter(ExamPaperID__in=exam_paper_obj_list, IsPassed=False)
        .values('ExamPaperID')
        .distinct()
    )
    to_finish_exam_paper_num = examinee_obj_list.count()
    # print(examinee_obj_list)
    context = {
        "to_exam_paper_num": to_exam_paper_num,
        "to_finish_exam_paper_num": to_finish_exam_paper_num,
    }
    # print(context)
    return JsonResponse(context)


def need_to_know(req):
    return render(req, "need_to_know.html")


# 培训


def train(request, bookid):
    context = {'bookid': bookid, 'book_min_time': 0}
    book_obj = Book.objects.filter(BookID=bookid).first()
    if book_obj is not None:
        context['book_min_time'] = book_obj.MinTime

    return render(request, "train.html", context)


def train_skill_category(request):
    return render(request, "train_skill_category.html")


def train_book(request):
    training_bank_id = request.GET.get('skill_type', 0)
    context = {'training_bank_id': training_bank_id}
    return render(request, "train_book.html", context)


def train_records(req):
    return render(req, "train_records.html")


# 训练


def practice(req):
    return render(req, "practice.html")


def practice_records(req):
    return render(req, "practice_records.html")


# 考核


def pre_issue_paper(request):
    return render(request, "pre_issue_paper.html")


def issue_paper(request):
    return render(request, "issue_paper.html")


def pre_exam(req):
    return render(req, "pre_exam.html")


def exam_notice(req):
    return render(req, "exam_notice.html")


def examing(req):
    return render(req, "examing.html")


def exam_records(req):
    return render(req, "exam_records.html")


def login(request):
    return render(request, "login.html")


def register(req):
    return render(req, "register.html")


# @login_required


def master(req):
    return render(req, "master.html")


def personal(req):
    return render(req, "personal.html")


def technology(req):
    return render(req, "technology.html")


def option(req):
    return render(req, "option.html")


def plant(req):
    return render(req, "plant.html")


@csrf_exempt
def questionbank(request):
    site = request.POST.get("Admin_Site")
    plant = request.POST.get("Admin_Plant")
    technologys = (
        Technology.objects.filter(Site=site, PlantName=plant)
        .values("Name")
        .order_by("Name")
    )
    Users = (
        User.objects.filter(Site=site, PlantName=plant)
        .values("StaffID", "StaffName")
        .order_by("StaffID", "StaffName")
    )
    return render(
        request, 'questionbank.html', {"technologys": technologys, "Users": Users}
    )


@csrf_exempt
def trainingbank(request):
    site = request.POST.get("Admin_Site")
    plant = request.POST.get("Admin_Plant")
    technologys = (
        Technology.objects.filter(Site=site, PlantName=plant)
        .values("Name")
        .order_by("Name")
    )
    users = (
        User.objects.filter(Site=site, PlantName=plant)
        .values("StaffID", "StaffName")
        .order_by("StaffID", "StaffName")
    )
    return render(
        request, 'trainingbank.html', {"technologys": technologys, "users": users}
    )


@csrf_exempt
def book(request, trainingBankID, trainingBankName):
    context = {'trainingBankID': trainingBankID, 'trainingBankName': trainingBankName}
    return render(request, "book.html", context)


@csrf_exempt
def examieed(request, ExamPaperID):
    # print('ExamPaperID:',ExamPaperID)
    context = {'ExamPaperID': ExamPaperID}
    return render(request, "examieed.html", context)


@csrf_exempt
def book_data(request, trainingBankID):
    is_training_material_maker = "false"
    staff_id = request.GET.get('Staff_ID', '')
    if staff_id != "":
        user_obj = User.objects.filter(StaffID=staff_id).first()
        if user_obj is not None and user_obj.Role == "Admin":
            is_training_material_maker = "true"
        else:
            training_bank_objs = TrainingBank.objects.filter(
                TrainingBankID=trainingBankID, TrainingMaterialMaker__contains=staff_id
            )
            if training_bank_objs.exists():
                is_training_material_maker = "true"

    # 获得点击的book的BookID,将查询结果返回给BookPage页面
    # 通过 id查询到bookPage
    books = Book.objects.filter(TrainingBankID=trainingBankID)
    result = []
    # 遍历结果
    for e in books:
        # 转成字典（json类型）
        jsongstr = model_to_dict(e)
        # 塞返回结果list
        result.append(jsongstr)
    # print(result)
    return JsonResponse(
        {'is_training_material_maker': is_training_material_maker, "book": result}
    )


def choiceQuestion(request, questionBankID, site, plant, name, isExercised):
    # 方法一 values()得到所有的值，不需要遍历
    # books = Book.objects.filter(Site=site,Plant=plant,TrainingBankName=name).values()
    books = Book.objects.filter(Site=site, Plant=plant, TrainingBankName=name).values()
    books = list(books)  # list类型
    # print("获取到教材：", books)
    # return JsonResponse(books, safe=False)
    """ 
    # 方法二 
    books = Book.objects.filter(TrainingBankName=name)
    # 遍历结果
    result = []
    for e in books:
        # 转成字典（json类型）
        jsongstr = model_to_dict(e)
        # 塞返回结果list
        result.append(jsongstr)
        print('result:%s' %(result))
    """
    context = {
        'questionBankID': questionBankID,
        'name': name,
        'isExercised': isExercised,
        'books': books,
    }
    return render(request, "choiceQuestion.html", context)


@csrf_exempt
def choiceQuestion_data(request, questionBankID):
    is_exam_question_maker = "false"
    staff_id = request.GET.get('Staff_ID', '')
    if staff_id != "":
        user_obj = User.objects.filter(StaffID=staff_id).first()
        if user_obj is not None and user_obj.Role == "Admin":
            is_exam_question_maker = "true"
        else:
            question_bank_objs = QuestionBank.objects.filter(
                QuestionBankID=questionBankID, ExamQuestionMaker__contains=staff_id
            )
            if question_bank_objs.exists():
                is_exam_question_maker = "true"

    test = ChoiceQuestion.objects.filter(QuestionBankID=questionBankID)
    result = []
    # 遍历结果
    for e in test:
        # 转成字典（json类型）
        jsongstr = model_to_dict(e)
        # 塞返回结果list
        result.append(jsongstr)
    # print(result)
    return JsonResponse(
        {"is_exam_question_maker": is_exam_question_maker, "test": result}
    )


def bookpage(request, bookid):
    context = {'bookid': bookid}
    book_obj = Book.objects.filter(BookID=bookid).first()
    if book_obj is not None:
        context['book'] = book_obj

    # // 这个函数现在点击哪一本书的链接就跳到bookPage页面，不做数据查询
    return render(request, "bookPage.html", context)


@csrf_exempt
def bookpage_data(request, bookid):
    is_training_material_maker = "false"
    staff_id = request.GET.get('Staff_ID', '')
    if staff_id != "":
        user_obj = User.objects.filter(StaffID=staff_id).first()
        if user_obj is not None and user_obj.Role == "Admin":
            is_training_material_maker = "true"
        else:
            training_bank_id = request.GET.get('training_bank_id', '')
            training_bank_objs = TrainingBank.objects.filter(
                TrainingBankID=training_bank_id,
                TrainingMaterialMaker__contains=staff_id,
            )
            if training_bank_objs.exists():
                is_training_material_maker = "true"

    # 获得点击的book的BookID,将查询结果返回给BookPage页面
    # 通过 id查询到bookPage
    bookPages = BookPage.objects.filter(BookID=bookid).order_by('Sequence')
    result = []
    # 遍历结果
    for e in bookPages:
        # 转成字典（json类型）
        # e.Content = settings.MEDIA_URL + e.Content
        # e.AudioPathFile = settings.MEDIA_URL + e.AudioPathFile
        jsongstr = model_to_dict(e)
        # 塞返回结果list
        result.append(jsongstr)
    # print(result)
    return JsonResponse(
        {"is_training_material_maker": is_training_material_maker, "bookPage": result}
    )


@csrf_exempt
def cqHistory(request, choiceQuestionID):
    context = {'choiceQuestionID': choiceQuestionID}
    cq_obj = ChoiceQuestion.objects.filter(ChoiceQuestionID=choiceQuestionID).first()
    if cq_obj is not None:
        context['ChoiceQuestion'] = cq_obj
    return render(request, "choiceHistory.html", context)


@csrf_exempt
def cqHistory_data(request, choiceQuestionID):
    cqHistory = ChoiceQuestionHistory.objects.filter(ChoiceQuestionID=choiceQuestionID)
    result = []
    # 遍历结果
    for e in cqHistory:
        # 转成字典（json类型）
        jsongstr = model_to_dict(e)
        # 塞返回结果list
        result.append(jsongstr)
    return JsonResponse({"cqHistory": result})


@csrf_exempt
def add_bookpage(request, bookid):
    if request.method == 'POST':
        # try:
        # return HttpResponse(repr(os.path.splitext(request.FILES['content'].name)))
        book_obj = Book.objects.get(pk=bookid)
        if book_obj is not None:
            training_book_path = "{0}/{1}/{2}".format(
                TRAINING_BANK, book_obj.TrainingBankName, book_obj.Name
            )
            absolute_training_book_path = "{0}/{1}".format(
                settings.MEDIA_ROOT, training_book_path
            )
            if not (os.path.exists(absolute_training_book_path)):
                os.makedirs(absolute_training_book_path)

            now = timezone.now()
            content_file_name = "{0}{1}".format(
                now.strftime('%Y%m%d%H%M%S%f'),
                os.path.splitext(request.FILES['content'].name)[-1],
            )
            content_file_path = "{0}/{1}".format(training_book_path, content_file_name)
            absolute_content_file_path = "{0}/{1}".format(
                settings.MEDIA_ROOT, content_file_path
            )
            with open(absolute_content_file_path, 'wb+') as destination:
                for chunk in request.FILES['content'].chunks():
                    destination.write(chunk)

            audio_file_path = ""
            if request.FILES.get('audio_file', None) is not None:
                audio_file_name = "{0}{1}".format(
                    now.strftime('%Y%m%d%H%M%S%f'),
                    os.path.splitext(request.FILES['audio_file'].name)[-1],
                )
                audio_file_path = "{0}/{1}".format(training_book_path, audio_file_name)
                absolute_audio_file_path = "{0}/{1}".format(
                    settings.MEDIA_ROOT, audio_file_path
                )
                with open(absolute_audio_file_path, 'wb+') as destination:
                    for chunk in request.FILES['audio_file'].chunks():
                        destination.write(chunk)

            sequence = request.POST['sequence']
            if request.POST['insert_direction'] == "up":
                BookPage.objects.filter(
                    BookID=bookid, Sequence__gte=int(sequence)
                ).update(Sequence=F('Sequence') + 1)
                insert_sequence = int(sequence)
            elif request.POST['insert_direction'] == "down":
                BookPage.objects.filter(
                    BookID=bookid, Sequence__gt=int(sequence)
                ).update(Sequence=F('Sequence') + 1)
                insert_sequence = int(sequence) + 1

            min_time = request.POST['mintime'] if request.POST['mintime'] != '' else 0
            bookpage_obj = BookPage.objects.create(
                BookID=bookid,
                Name=request.POST['name'],
                Content=content_file_path,
                AudioPathFile=audio_file_path,
                MinTime=min_time,
                Sequence=insert_sequence,
                CreatorID=request.POST['creator_id'],
                CreatorName=request.POST['creator_name'],
                UpdateDate=timezone.now(),
                CreateDate=timezone.now(),
            )

            if bookpage_obj is not None:
                bookpage_history_obj = BookPageHistory.objects.create(
                    BookID=bookid,
                    BookPageID=bookpage_obj.BookPageID,
                    Reason=request.POST['reason'],
                    Operation="Added",
                    StaffID=request.POST['creator_id'],
                    StaffName=request.POST['creator_name'],
                    CreateDate=timezone.now(),
                )

        return JsonResponse({"result": "true"})
        # except Exception as e:
        # return HttpResponse(repr(e))
    else:
        return JsonResponse({"result": "false"})


def book_oprating_record(request, bookid):
    context = {'bookid': bookid}
    book_obj = Book.objects.filter(BookID=bookid).first()
    if book_obj is not None:
        context['book'] = book_obj
    return render(request, 'bookPageHistory.html', context)


@csrf_exempt
def choiceHistory(request, questionBankID, name, isExercised):
    context = {
        'questionBankID': questionBankID,
        'name': name,
        'isExercised': isExercised,
    }
    choiceQuestion_obj = ChoiceQuestion.objects.filter(
        QuestionBankID=questionBankID
    ).values("ChoiceQuestionID")

    if choiceQuestion_obj is not None:
        context = {
            'questionBankID': questionBankID,
            'name': name,
            'isExercised': isExercised,
        }
    return render(request, 'qBankHistory.html', context)


@csrf_exempt
def choiceHistory_data(request, questionBankID, name, isExercised):
    try:
        choiceQuestion_obj = ChoiceQuestion.objects.filter(
            QuestionBankID=questionBankID
        ).values_list('ChoiceQuestionID')
        result = []
        for choiceQuestionID in choiceQuestion_obj:
            cHinstory_list = ChoiceQuestionHistory.objects.filter(
                ChoiceQuestionID=choiceQuestionID[0]
            )
            # 遍历结果
            for e in cHinstory_list:
                # 转成字典（json类型）
                jsongstr = model_to_dict(e)
                # 塞返回结果list
                result.append(jsongstr)
        return JsonResponse({"cHinstory_list": result})
    except Exception as e:
        return HttpResponse(repr(e))


@csrf_exempt
def book_oprating_record_data(request, bookid):
    # 获取指定book和bookpage的history
    try:
        result = []
        book_history_list = BookHistory.objects.filter(BookID=bookid).order_by(
            "CreateDate"
        )
        # 遍历结果
        for e in book_history_list:
            # 转成字典（json类型）
            jsongstr = model_to_dict(e)
            # 塞返回结果list
            result.append(jsongstr)
        # print(result)
        book_page_history_list = BookPageHistory.objects.filter(BookID=bookid).order_by(
            "CreateDate"
        )
        # 遍历结果
        for e in book_page_history_list:
            # 转成字典（json类型）
            jsongstr = model_to_dict(e)
            # 塞返回结果list
            result.append(jsongstr)
        return JsonResponse({"data": result})
    except Exception as e:
        return HttpResponse(repr(e))


@csrf_exempt
def skill_sort(request):
    return render(request, "skill_sort.html")


@csrf_exempt
def qBankSort(request):
    return render(request, "qBankSort.html")


@csrf_exempt
def bookPage(request):
    # 这个函数现在点击哪一本书的链接就跳到bookPage页面，不做数据查询
    return render(request, "bookPage.html")


# 查询所有技能教材页


@csrf_exempt
def get_all_bookpage(request):
    try:
        bookpageList = serializers.serialize("json", BookPage.objects.all())
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": bookpageList})
    # return render(req, "book.html")


@csrf_exempt
def book_exist(request):
    if request.method != 'POST':
        return JsonResponse({"result": ""})

    book_file_name = request.POST['book_file_name']
    staff_id = request.POST.get('Staff_ID', '')
    admin_site = request.POST.get('admin_site')
    admin_plant = request.POST.get('admin_plant')
    book_obj = Book.objects.filter(
        Site=admin_site, Plant=admin_plant, Name=book_file_name
    ).first()
    if book_obj is not None:
        training_bank_obj_list = TrainingBank.objects.filter(
            TrainingBankID=book_obj.TrainingBankID,
            TrainingMaterialMaker__contains=staff_id,
        )
        is_book_maker = training_bank_obj_list.exists()
        return JsonResponse(
            {
                "result": "true",
                'is_book_maker': is_book_maker,
                "training_bank_name": book_obj.TrainingBankName,
            }
        )
    else:
        return JsonResponse({"result": "false", 'is_book_maker': False})


# 新增技能教材add_book


@csrf_exempt
def add_book(request):
    # request.encoding = 'utf-8'
    if request.method == 'POST':
        # 数据库存文件名，文件名 = inputName.name，inputName是传入的一整个文件

        reason = request.POST['reason']
        book_mintime = request.POST['book_mintime']
        creator_id = request.POST['creator_id']
        creator_name = request.POST['creator_name']
        creator_site = request.POST['site']
        creator_plant = request.POST['plant']

        training_bank_id = request.POST['training_bank_id']
        book_zip_file = request.FILES.get('book_zip_file')

        training_bank_obj = TrainingBank.objects.filter(
            TrainingBankID=training_bank_id
        ).first()
        if training_bank_obj is None:
            result = {"code": 1, "text": "教材库不存在!"}
            return JsonResponse({"result": result})

        training_bank_path = "{0}/{1}".format(TRAINING_BANK, training_bank_obj.Name)
        absolute_training_bank_path = "{0}/{1}".format(
            settings.MEDIA_ROOT, training_bank_path
        )
        if not (os.path.exists(absolute_training_bank_path)):
            os.makedirs(absolute_training_bank_path)

        book_zip_file_path = "{0}/{1}".format(
            absolute_training_bank_path, book_zip_file.name
        )
        with open(book_zip_file_path, 'wb+') as destination:
            for chunk in book_zip_file.chunks():
                destination.write(chunk)

        book_name = os.path.splitext(book_zip_file.name)[0]
        training_book_path = "{0}/{1}".format(training_bank_path, book_name)
        absolute_training_book_path = "{0}/{1}/".format(
            absolute_training_bank_path, book_name
        )
        book_obj = Book.objects.filter(
            Site=creator_site, Plant=creator_plant, Name=book_name
        ).first()
        if book_obj is not None:
            Book.objects.filter(
                Site=creator_site, Plant=creator_plant, Name=book_name
            ).update(
                TrainingBankID=training_bank_id,
                TrainingBankName=training_bank_obj.Name,
                MinTime=book_mintime,
                CreatorName=creator_name,
                CreatorID=creator_id,
                UpdateDate=timezone.now(),
            )
            book_page_obj_list = BookPage.objects.filter(BookID=book_obj.BookID)
            book_page_obj_list.delete()
            shutil.rmtree(absolute_training_book_path, True)  # 最后删除总文件夹

        image_file_extension = ['.png', '.jpg']
        audio_file_extension = ['.m4a', '.wav']
        book_page_content_dict = {}
        book_page_audio_dict = {}

        # 解压压缩文件
        book_zip_file_obj = zipfile.ZipFile(book_zip_file_path)
        book_dir = book_name + "/"
        for name in book_zip_file_obj.namelist():
            if book_dir in name:
                file_name = os.path.basename(name)
                file_info_list = os.path.splitext(file_name)

                sequence_list = re.findall(u'^投影片\d+|p\d+$', file_info_list[0], re.I)
                if sequence_list:
                    sequence = re.sub(u'^投影片|p', "", sequence_list[0], 0, re.I)
                    # print(sequence)

                    if re.findall(r'^\.png|\.jpg$', file_info_list[1], re.I):
                        book_page_content_dict[sequence] = file_info_list
                    elif re.findall(r'^\.m4a|\.wav$', file_info_list[1], re.I):
                        book_page_audio_dict[sequence] = file_info_list

                book_zip_file_obj.extract(name, absolute_training_bank_path)

        # print(book_page_content_dict)
        # print(book_page_audio_dict)

        book_zip_file_obj.close()
        # 删除压缩文件
        if os.path.exists(book_zip_file_path):
            os.remove(book_zip_file_path)

        if len(book_page_content_dict) == 0:
            shutil.rmtree(absolute_training_book_path, True)  # 最后删除总文件夹
            text = '教材"%s"没有教材页，无法添加!' % (book_name)
            result = {"code": 2, "text": text}
            return JsonResponse({"result": result})

        if book_obj is None:
            book_obj = Book.objects.create(
                TrainingBankID=training_bank_id,
                TrainingBankName=training_bank_obj.Name,
                Site=creator_site,
                Plant=creator_plant,
                Name=book_name,
                MinTime=book_mintime,
                CreatorName=creator_name,
                CreatorID=creator_id,
                UpdateDate=timezone.now(),
                CreateDate=timezone.now(),
            )

        if book_obj is not None:
            book_history_obj = BookHistory.objects.create(
                BookID=book_obj.BookID,
                Reason=reason,
                Operation="Added",
                StaffID=creator_id,
                StaffName=creator_name,
                CreateDate=timezone.now(),
            )

            for sequence, value in book_page_content_dict.items():
                time.sleep(0.001)
                now = timezone.now()
                content_file_name = "{0}{1}".format(
                    now.strftime('%Y%m%d%H%M%S%f'), value[1]
                )

                content_file_path = "{0}/{1}".format(
                    training_book_path, content_file_name
                )
                # print(content_file_path)

                os.rename(
                    absolute_training_book_path + ''.join(value),
                    absolute_training_book_path + content_file_name,
                )

                book_page_audio = book_page_audio_dict.get(sequence)
                if book_page_audio is not None:
                    audio_file_name = "{0}{1}".format(
                        now.strftime('%Y%m%d%H%M%S%f'), book_page_audio[1]
                    )
                    audio_file_path = "{0}/{1}".format(
                        training_book_path, audio_file_name
                    )

                    os.rename(
                        absolute_training_book_path + ''.join(book_page_audio),
                        absolute_training_book_path + audio_file_name,
                    )
                else:
                    audio_file_path = ""
                # print(audio_file_path)

                bookpage_obj = BookPage.objects.create(
                    BookID=book_obj.BookID,
                    Name="",
                    Content=content_file_path,
                    AudioPathFile=audio_file_path,
                    MinTime=0,
                    Sequence=sequence,
                    CreatorID=creator_id,
                    CreatorName=creator_name,
                    UpdateDate=timezone.now(),
                    CreateDate=timezone.now(),
                )

                if bookpage_obj is not None:
                    bookpage_history_obj = BookPageHistory.objects.create(
                        BookID=book_obj.BookID,
                        BookPageID=bookpage_obj.BookPageID,
                        Reason="Added {0} book page".format(book_name),
                        Operation="Added",
                        StaffID=creator_id,
                        StaffName=creator_name,
                        CreateDate=timezone.now(),
                    )

        # return JsonResponse({"result":"true"})
        result = {"code": 0, "text": "添加成功!"}
        return JsonResponse({"result": result})
        # return HttpResponse({"data":"ok"})
    # else:
    #     result = {"code":1, "text": "添加不成功!"}
    #     return JsonResponse({"result": result})
    # return JsonResponse({"result":"false"})


# # 检查数据库是否存在这个文件
# @csrf_exempt
# def choiceQuestion_exist(request):
#     try:
#         if request.method != 'POST':
#             return JsonResponse({"result": ""})
#         ppt_name = request.POST['ppt_name']
#         choiceQuestions = ChoiceQuestion.objects.all().values_list('Content')
#         res =''
#         # 遍历结果
#         for e in choiceQuestions:
#             i = str(e)
#             pname = '/'.join(i.split('/')[-2:-1])    #截取地址中间的文件名部分
#             # print('i:',i)
#             # print('pname:',pname)
#             if  ppt_name == pname:
#                 res = 'true'
#                 break
#             else:
#                 res = 'false'
#         return JsonResponse({"result": res})
#     except Exception as e:
#         print('出错了: %s' % e)
#         return HttpResponse("1:" + " " + repr(e))


# 新增技能考题
@csrf_exempt
def add_choiceQuestion(request):
    # try:
    request.encoding = 'utf-8'
    inputQuestionBankID = request.POST['newQuestionBankID']
    inputBookID = request.POST['newBookID']
    inputSubject = request.POST['newSubject']
    inputType = request.POST['newType']
    # 数据库存文件名，文件名 = inputName.name，inputName是传入的一整个文件
    inputContent = request.FILES.get('newContent')
    fileName = inputContent.name
    # fname = os.path.splitext(fileName)[0]
    inputCreatorName = request.POST['newCreatorName']
    inputCreatorID = request.POST['newCreatorID']
    inputReason = request.POST['newReason']
    inputCreateDate = timezone.now()
    inputUpdateDate = timezone.now()
    isExercise = ''

    # 判断是否存在该题库
    question_bank_obj = QuestionBank.objects.filter(
        QuestionBankID=inputQuestionBankID
    ).first()
    book_obj = Book.objects.get(BookID=inputBookID)
    result = question_bank_obj.IsExercised
    if question_bank_obj is None:
        return JsonResponse(
            {
                "status": -1,
                "data": "",
                "message": "该数据库不存在!\nThe database does not exist!",
            }
        )
    # 判断是该题库的训练题库还是考核题库
    if result:
        isExercise = 'exercise'
    else:
        isExercise = 'exam'

    # 题库地址：bank/AXI/exam/BookName/jpg
    # settings.MEDIA_ROOT 是E:\ATQS\trunk\sourcecode\itqs\itqs\media
    question_bank_path = "{0}/{1}/{2}".format(
        QUESTION_BANK, question_bank_obj.Name, isExercise
    )
    absolute_question_bank_path = "{0}/{1}".format(
        settings.MEDIA_ROOT, question_bank_path
    )
    if not (os.path.exists(absolute_question_bank_path)):
        os.makedirs(absolute_question_bank_path)
    # question_bank/SPI
    # E:\ATQS\trunk\sourcecode\itqs\itqs\media/question_bank/SPI

    # 创建choiceQuestion目录,并创建该ppt读取到的图片存放位置
    question_choiceQuestion_path = "{0}/{1}/".format(question_bank_path, book_obj.Name)
    absolute_question_choiceQuestion_path = "{0}/{1}/".format(
        absolute_question_bank_path, book_obj.Name
    )
    if not (os.path.exists(absolute_question_choiceQuestion_path)):
        os.makedirs(absolute_question_choiceQuestion_path)

    # 创建pptx目录，写入该文件夹
    pptx_file_path = "{0}/{1}".format(absolute_question_choiceQuestion_path, fileName)
    with open(pptx_file_path, 'wb+') as destination:
        for chunk in inputContent.chunks():
            destination.write(chunk)

    optionID = 0
    i = 0
    prs = Presentation(pptx_file_path)
    question_list = []
    for slide in prs.slides:
        option_list = []
        urls = ""
        i += 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                None
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                blob = image.blob
                picture_name = '%s_%s_%s.%s' % (
                    slide.slide_id,
                    shape.shape_id,
                    shape.name,
                    image.ext,
                )
                suffix = os.path.splitext(picture_name)[1]
                time.sleep(0.001)
                now = timezone.now()
                file_name = "{0}{1}".format(now.strftime('%Y%m%d%H%M%S%f'), suffix)

                # 图片写入url地址的文件夹
                urls = question_choiceQuestion_path + file_name
                absolute_urls = absolute_question_choiceQuestion_path + file_name
                with open(absolute_urls, 'wb+') as file:
                    file.write(blob)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for cell in shape.table.iter_cells():
                    option_list.append(cell.text)
        if urls != "" and option_list and "Y" in ''.join(option_list):
            question_dict = {'urls': urls, 'option_list': option_list}
            question_list.append(question_dict)
        else:
            if urls == '' and len(option_list) != 0:
                text = (
                    '第%s张幻灯片无法读取图片！请修改后重新上传!\nSlide %s is missing picture! Please check and upload again!'
                    % (i, i)
                )
            elif urls != '' and len(option_list) == 0:
                text = (
                    '第%s张幻灯片无法读取选项！请修改后重新上传!\nSlide %s is missing options! Please check and upload again!'
                    % (i, i)
                )
            elif urls == '' and len(option_list) == 0:
                text = (
                    '第%s张幻灯片无法读取选项和图片！请修改后重新上传!\nSlide %s is missing options and pictures! Please check and upload again!'
                    % (i, i)
                )
            elif urls != "" and len(option_list) and "Y" not in ''.join(option_list):
                text = (
                    '第%s张幻灯片的选项没有答案！请修改后重新上传!There are no answers to the options in slide %s! Please check and upload again!'
                    % (i, i)
                )
            return JsonResponse({"status": -1, "data": "", "message": text})
    for question in question_list:
        # 保存存入的数据到表ChoiceQuestion
        choiceQuestion = ChoiceQuestion.objects.create(
            QuestionBankID=inputQuestionBankID,
            BookID=inputBookID,
            Subject=inputSubject,
            Type=inputType,
            Content=question['urls'],
            CreatorName=inputCreatorName,
            CreatorID=inputCreatorID,
            UpdateDate=inputUpdateDate,
            CreateDate=inputCreateDate,
        )
        optionID = choiceQuestion.ChoiceQuestionID

        # 保存存入的数据到表QuestionOption
        for option in question['option_list']:
            if "Y" in option:
                isAnswer = "True"
                op_list = option.split(']')[1:]
                op_char = ''.join(op_list)
                op = op_char.strip(' ')
            else:
                isAnswer = "False"
                op_list = option.split(']')[1:]
                op_char = ''.join(op_list)
                op = op_char.strip(' ')
            questionOption = QuestionOption.objects.create(
                ChoiceQuestionID=optionID, Option=op, IsAnswer=isAnswer
            )

        # 将理由存入表ChoiceQuestionHistory
        choiceQuestionHistory = ChoiceQuestionHistory.objects.create(
            ChoiceQuestionID=optionID,
            Reason=inputReason,
            Operation='Added',
            StaffName=inputCreatorName,
            StaffID=inputCreatorID,
            CreateDate=inputCreateDate,
        )

    # 读取.pptx文件的内容后删除
    if os.path.exists(pptx_file_path):
        os.remove(pptx_file_path)

    # except Exception as e:
    #     print('出错了: %s' % e)
    #     return HttpResponse("1:" + " " + repr(e))
    return JsonResponse({"status": 0, "data": "ok", "message": ""})


# # 修改技能考题
# @csrf_exempt
# def modify_choiceQuestion(request):
#     request.encoding = 'utf-8'
#     # book_id = request.POST.get('modifyBookID')
#     choiceQuestionID = request.POST.get('choiceQuestionID')
#     bookID = request.POST.get('modifyBID')

#     creator_id = request.POST.get('modifyCreatorID')
#     creator_name = request.POST.get('modifyCreatorName')
#     modify_reason = request.POST.get('modifyreason')

#     choiceQuestion_obj = ChoiceQuestion.objects.filter(ChoiceQuestionID=choiceQuestionID).update(
#         BookID=bookID, UpdateDate=timezone.now())

#     # choiceQuestion = ChoiceQuestion.objects.filter(ChoiceQuestionID=ChoiceQuestionID).update(
#     #     Subject=Subject, Type=Type, BookID=BookID, UpdateDate=timezone.now())
#     if choiceQuestion_obj is not None:
#         choiceQuestionHistory = ChoiceQuestionHistory.objects.create(
#             ChoiceQuestionID=choiceQuestionID,
#             Reason=modify_reason,
#             Operation="modified",
#             StaffID=creator_id,
#             StaffName=creator_name,
#             CreateDate=timezone.now()
#         )

#     return JsonResponse({"result":"true"})

# 获得该幻灯片的所有选项


@csrf_exempt
def get_options(req):
    req.encoding = 'utf-8'
    try:
        ChoiceQuestionID = req.GET['ChoiceQuestionID']
        options = QuestionOption.objects.filter(ChoiceQuestionID=ChoiceQuestionID)
        result = []
        # 遍历结果
        for e in options:
            # 转成字典（json类型）
            jsongstr = model_to_dict(e)
            # 塞返回结果list
            result.append(jsongstr)
        return JsonResponse({"options": result})
    except Exception as e:
        return HttpResponse(repr(e))


# 获得适用书籍的名字


@csrf_exempt
def get_bookName(req):
    req.encoding = 'utf-8'
    try:
        BookID = req.GET['BookID']
        book = Book.objects.filter(BookID=BookID)
        result = []
        # 遍历结果
        for e in book:
            # 转成字典（json类型）
            jsongstr = model_to_dict(e)
            # 塞返回结果list
            result.append(jsongstr)
        return JsonResponse({"book": result})
    except Exception as e:
        return HttpResponse(repr(e))


# 查询所有技能考题历史记录


@csrf_exempt
def get_choiceHistory(req):
    req.encoding = 'utf-8'
    try:
        # 获取全部的历史纪录
        choiceHistory = ChoiceQuestionHistory.objects.filter()
        result = []
        # 遍历结果
        for e in choiceHistory:
            # 转成字典（json类型）
            jsongstr = model_to_dict(e)
            # 塞返回结果list
            result.append(jsongstr)
        # print(result)
        return JsonResponse({"choiceHistory": result})
    except Exception as e:
        return HttpResponse(repr(e))


@csrf_exempt
def book_question_num(request):
    try:
        book_id = request.POST.get('book_id', 0)
        question_bank_id = request.POST.get('question_bank_id', '')
        choice_question_obj_list = ChoiceQuestion.objects.filter(
            QuestionBankID=question_bank_id, BookID=book_id
        )

        question_option_obj_list = (
            QuestionOption.objects.filter(
                ChoiceQuestionID__in=choice_question_obj_list, IsAnswer=True
            )
            .values('Option')
            .annotate(Count=Count('QuestionOptionID'))
            .values('Option', 'Count')
        )

        # print(question_option_obj_list)
        return JsonResponse(
            {
                "book_question_num": choice_question_obj_list.count(),
                "question_option_count_list": list(question_option_obj_list),
            }
        )
    except Exception as e:
        return HttpResponse(repr(e))


# 删除技能教材(book) debug mwj


@csrf_exempt
def del_book(request):
    request.encoding = 'utf-8'
    # 获得被删除的书本ID
    BookID = request.POST.getlist('BookID[]', [])
    reason = request.POST.get('reason', '')
    creator_id = request.POST.get('creator_id', '')
    creator_name = request.POST.get('creator_name', '')
    errMsgListt = []

    for bookid in BookID:
        book_obj = Book.objects.filter(BookID=bookid).first()
        if book_obj is None:
            return JsonResponse(
                {
                    "status": -1,
                    "data": "false",
                    "message": "删除失败,有教材已经不存在，请刷新后重新选择!\nFail to delete. Some teaching materials have been deleted. Please refresh and select again!",
                }
            )

        choice_question_id = ChoiceQuestion.objects.filter(BookID=bookid).values_list(
            'QuestionBankID'
        )
        if len(choice_question_id) > 0:
            choice_question_id = choice_question_id.distinct()
            # 在QuestionBank存在，且是用于考核的
            axam_bank = QuestionBank.objects.filter(
                QuestionBankID__in=choice_question_id, IsExercised=False
            ).values()
            exercise_bank = QuestionBank.objects.filter(
                QuestionBankID__in=choice_question_id, IsExercised=True
            ).values()
            axam_bank = list(axam_bank)
            exercise_bank = list(exercise_bank)

            axam_bank_len = len(axam_bank)
            exercise_bank_len = len(exercise_bank)
            if axam_bank_len > 0 and exercise_bank_len > 0:
                text = (
                    '教材"%s",已被"%s考核库"和"%s训练库"使用，不能删除! \n "%s" is used in "%s question bank (exam)" and "%s question bank (practice)", and cannot be deleted!'
                    % (
                        book_obj.Name,
                        axam_bank[0]["Name"],
                        exercise_bank[0]["Name"],
                        book_obj.Name,
                        axam_bank[0]["Name"],
                        exercise_bank[0]["Name"],
                    )
                )
                errMsgListt.append(text)
                continue
            elif axam_bank_len > 0 and exercise_bank_len == 0:
                text = (
                    '教材"%s",已被"%s考核库"使用，不能删除！\n "%s" is used in "%s question bank (exam)", and cannot be deleted!'
                    % (
                        book_obj.Name,
                        axam_bank[0]["Name"],
                        book_obj.Name,
                        axam_bank[0]["Name"],
                    )
                )
                errMsgListt.append(text)
                continue
            elif axam_bank_len == 0 and exercise_bank_len > 0:
                text = (
                    '教材"%s",已被"%s训练库"使用，不能删除！\n "%s" is used in "%s question bank (practice)" and cannot be deleted!'
                    % (
                        book_obj.Name,
                        exercise_bank[0]["Name"],
                        book_obj.Name,
                        exercise_bank[0]["Name"],
                    )
                )
                errMsgListt.append(text)
                continue

        else:
            book_name = book_obj.Name
            training_bank_name = book_obj.TrainingBankName

            absolute_training_book_path = "{0}/{1}/{2}/{3}".format(
                settings.MEDIA_ROOT, TRAINING_BANK, training_bank_name, book_name
            )
            shutil.rmtree(absolute_training_book_path, True)  # 最后删除总文件夹

            bookpage_obj_list = BookPage.objects.filter(BookID=bookid)
            # 删除该本书的page
            bookpage_obj_list.delete()

            book_history_obj = BookHistory.objects.create(
                BookID=book_obj.BookID,
                Reason=reason,
                Operation="Deleted",
                StaffID=creator_id,
                StaffName=creator_name,
                CreateDate=timezone.now(),
            )
            book_obj.delete()

    if len(errMsgListt) > 0:
        return JsonResponse({"status": -1, "data": "error", "message": errMsgListt})
    return JsonResponse({"status": 0, "data": "ok", "message": ""})


# 删除技能考题(ChoiceQuestion)
@csrf_exempt
def del_choiceQuestion(req):
    # print(req.POST)    ##这个可以看到req里面都有什么
    req.encoding = 'utf-8'
    if req.method != 'POST':
        return JsonResponse({"status": -1, "data": "false", "message": ""})
    # 获得被删除的考题ID
    ChoiceQuestionID = req.POST.getlist('ChoiceQuestionID[]', [])
    Reason = req.POST.get('delReason')
    creator_id = req.POST.get('creator_id')
    creator_name = req.POST.get('creator_name')

    for delID in ChoiceQuestionID:
        choiceQuestion = ChoiceQuestion.objects.filter(ChoiceQuestionID=delID)
        if len(choiceQuestion) == 0:
            return JsonResponse(
                {
                    "status": -1,
                    "data": "false",
                    "message": "你选的考题中有的考题已经不存在，请刷新后重新选择!\nSome of the questions you selected have been deleted. Please refresh and select again!",
                }
            )
        choiceQuestion_obj = choiceQuestion.values()
        choiceQuestion_obj = list(choiceQuestion_obj)
        questionOption_obj = QuestionOption.objects.filter(ChoiceQuestionID=delID)

        # 记录删除的操作
        choiceQuestionHistory_obj = ChoiceQuestionHistory.objects.create(
            ChoiceQuestionID=delID,
            Reason=Reason,
            Operation="Deleted",
            StaffID=creator_id,
            StaffName=creator_name,
            CreateDate=timezone.now(),
        )

        # 如果该题有在用就不能删除文件，Tab ExamedChoiceQuestion要使用该文件
        examedChoiceQuestion_obj = ExamedChoiceQuestion.objects.filter(
            ChoiceQuestionID=delID
        ).first()
        if examedChoiceQuestion_obj is None:
            if choiceQuestion_obj[0] is not None:
                path = "{0}/{1}".format(
                    settings.MEDIA_ROOT, choiceQuestion_obj[0]["Content"]
                )
                work_path = '/'.join(path.split('/')[:-1])
                if os.path.exists(path):  # 如果文件存在
                    os.remove(path)
                    if not os.listdir(work_path):
                        # os.remove(work_path)      # 如果文件夹为空，则删除文件夹
                        shutil.rmtree(work_path, True)
                    else:
                        print('The folder is not empty：%s' % work_path)
                else:
                    print('no such file:%s' % path)  # 则返回文件不存在

        # 删除数据库的这一条数据
        choiceQuestion.delete()
        questionOption_obj.delete()

    return JsonResponse({"status": 0, "data": "ok", "message": ""})


# 根据book删除技能考题
@csrf_exempt
def del_choiceQuestion_by_book(req):
    req.encoding = 'utf-8'
    if req.method != 'POST':
        return JsonResponse({"status": -1, "data": "false", "message": ""})
    bookid = req.POST.get('bookid')
    Reason = req.POST.get('delreason')
    creator_id = req.POST.get('creator_id')
    creator_name = req.POST.get('creator_name')
    ChoiceQuestionID = []
    allID = ChoiceQuestion.objects.filter(BookID=bookid).values()
    for i in allID:
        ChoiceQuestionID.append(i['ChoiceQuestionID'])
    for delID in ChoiceQuestionID:
        choiceQuestion = ChoiceQuestion.objects.filter(ChoiceQuestionID=delID)
        if len(choiceQuestion) == 0:
            return JsonResponse(
                {
                    "status": -1,
                    "data": "false",
                    "message": "你选的考题中有的考题已经不存在，请刷新后重新选择!\nSome of the questions you selected have been deleted. Please refresh and select again!",
                }
            )
        choiceQuestion_obj = choiceQuestion.values()
        choiceQuestion_obj = list(choiceQuestion_obj)
        questionOption_obj = QuestionOption.objects.filter(ChoiceQuestionID=delID)

        # 记录删除的操作
        choiceQuestionHistory_obj = ChoiceQuestionHistory.objects.create(
            ChoiceQuestionID=delID,
            Reason=Reason,
            Operation="Deleted",
            StaffID=creator_id,
            StaffName=creator_name,
            CreateDate=timezone.now(),
        )

        # 如果该题有在用就不能删除文件，Tab ExamedChoiceQuestion要使用该文件
        examedChoiceQuestion_obj = ExamedChoiceQuestion.objects.filter(
            ChoiceQuestionID=delID
        ).first()
        if examedChoiceQuestion_obj is None:
            if choiceQuestion_obj[0] is not None:
                path = "{0}/{1}".format(
                    settings.MEDIA_ROOT, choiceQuestion_obj[0]["Content"]
                )
                work_path = '/'.join(path.split('/')[:-1])
                if os.path.exists(path):  # 如果文件存在
                    os.remove(path)
                    if not os.listdir(work_path):
                        # os.remove(work_path)      # 如果文件夹为空，则删除文件夹
                        shutil.rmtree(work_path, True)
                    else:
                        print('The folder is not empty：%s' % work_path)
                else:
                    print('no such file:%s' % path)  # 则返回文件不存在

        # 删除数据库的这一条数据
        choiceQuestion.delete()
        questionOption_obj.delete()

    return JsonResponse({"status": 0, "data": "ok", "message": ""})


# 删除技能教材页(bookPage)
@csrf_exempt
def del_BookPage(request):
    request.encoding = 'utf-8'
    if request.method == 'POST':
        # 获得被删除的书本ID
        bookpage_id = request.POST['bookpage_id']
        bookpage_obj = BookPage.objects.get(BookPageID=bookpage_id)
        if bookpage_obj is not None:
            absolute_content_file_path = "{0}/{1}".format(
                settings.MEDIA_ROOT, bookpage_obj.Content
            )
            if os.path.exists(absolute_content_file_path):  # 如果文件存在
                os.remove(absolute_content_file_path)
            if bookpage_obj.AudioPathFile != "":
                absolute_audio_file_path = "{0}/{1}".format(
                    settings.MEDIA_ROOT, bookpage_obj.AudioPathFile
                )
                if os.path.exists(absolute_audio_file_path):  # 如果文件存在
                    os.remove(absolute_audio_file_path)

            bookpage_history_obj = BookPageHistory.objects.create(
                BookID=bookpage_obj.BookID,
                BookPageID=bookpage_obj.BookPageID,
                Reason=request.POST['reason'],
                Operation="Deleted",
                StaffID=request.POST['creator_id'],
                StaffName=request.POST['creator_name'],
                CreateDate=timezone.now(),
            )

            bookpage_obj.delete()

        return JsonResponse({"result": "true"})
    else:
        return JsonResponse({"result": "false"})


# 修改技能教材


@csrf_exempt
def modify_book(request):
    request.encoding = 'utf-8'

    book_id = request.POST.get('modifyBookID')
    book_min_time = request.POST.get("modifyBookMinTime")
    reason = request.POST.get('modifyreason')
    creator_id = request.POST.get('modifyCreatorID')
    creator_name = request.POST.get('modifyCreatorName')

    book_obj = Book.objects.filter(BookID=book_id).update(
        MinTime=book_min_time, UpdateDate=timezone.now()
    )
    if book_obj is not None:
        book_history_obj = BookHistory.objects.create(
            BookID=book_id,
            Reason=reason,
            Operation="Modify",
            StaffID=creator_id,
            StaffName=creator_name,
            CreateDate=timezone.now(),
        )

    return JsonResponse({"result": "true"})


# 修改技能教材页
@csrf_exempt
def modify_bookPage(req):
    req.encoding = 'utf-8'
    # 先通过ID查到所要修改的书本，修改书内的其他属性
    inputBookPageID = req.POST['modifyBookPageID']
    inputPageAudio = req.POST['modifyPageAudio']
    inputPageMinTime = req.POST['modifyPageMinTime']
    inputPageMinTime = req.POST['modifyPageMinTime']
    inputUpdateDate = timezone.now()
    # 查询出来需要需改的bookPageID
    modifyBookPage = BookPage.objects.get(BookPageID=inputBookPageID)
    # 获得修改后的数据
    modifyBookPage.AudioPathFile = inputPageAudio
    modifyBookPage.MinTime = inputPageMinTime
    modifyBookPage.UpdateDate = inputUpdateDate
    modifyBookPage.save()
    return JsonResponse({"data": "ok"})


# 查询所有技能类型


@csrf_exempt
def get_train_skill_type(request):
    site = request.POST.get("site")
    plant = request.POST.get("plant")
    skill_type = TrainingBank.objects.filter(
        Site=site, Plant=plant
    ).values()  # QuerySet类型
    return JsonResponse(list(skill_type), safe=False)


# 培训，根据技能类型查找教材


@csrf_exempt
def get_train_book(request):
    training_bank_id = request.POST.get('training_bank_id', 0)
    book_obj_list = Book.objects.filter(
        TrainingBankID=training_bank_id
    ).values()  # QuerySet类型
    book_list = list(book_obj_list)  # list类型
    return JsonResponse(book_list, safe=False)


# 培训，根据教材id查找教材pages


@csrf_exempt
def get_train_pages(request):
    bookID = request.POST.get('bookID', 0)
    pages = (
        BookPage.objects.filter(BookID=bookID).values().order_by("Sequence")
    )  # QuerySet类型
    pages = list(pages)
    # print('get_train_pages')
    print(pages)
    return JsonResponse(pages, safe=False)


# 培训，添加培训记录


@csrf_exempt
def add_train_record(request):
    if request.method == "POST":
        bookID = request.POST['bookID']
        staff_ID = request.POST['staff_ID']
        staff_name = request.POST['staff_name']
        train_time = request.POST['train_time']
        is_pass = request.POST['is_pass']
        book_obj = Book.objects.filter(BookID=bookID).first()
        if book_obj is not None:
            training_record_obj = TrainingRecord.objects.create(
                Site=book_obj.Site,
                Plant=book_obj.Plant,
                BookID=book_obj.BookID,
                BookName=book_obj.Name,
                BookMinTime=book_obj.MinTime,
                TrainingBankName=book_obj.TrainingBankName,
                StaffID=staff_ID,
                StaffName=staff_name,
                TrainingTime=train_time,
                IsPassed=is_pass,
                CreateDate=timezone.now(),
            )
            return JsonResponse({"result": "true"})
        else:
            return JsonResponse({"result": "false"})
    else:
        return JsonResponse({"result": ""})


# 培训，查找培训记录
@csrf_exempt
def get_train_record(request):
    role = request.POST.get('role')
    admin_site = request.POST.get('admin_site')
    admin_plant = request.POST.get('admin_plant')
    staff_ID = request.POST.get('staff_ID')
    try:
        if role == 'Admin':
            all_train_record = serializers.serialize(
                "json",
                TrainingRecord.objects.filter(
                    Site=admin_site, Plant=admin_plant
                ).order_by('-CreateDate'),
            )
        else:
            all_train_record = serializers.serialize(
                "json",
                TrainingRecord.objects.filter(StaffID=staff_ID).order_by('-CreateDate'),
            )
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": all_train_record})


def get_user_role(staff_id):
    if staff_id is None:
        return None

    user_obj = User.objects.filter(StaffID=staff_id).first()
    if user_obj is None:
        return None

    return user_obj.Role


# 训练，获取所有用于训练的book


@csrf_exempt
def get_all_practice_books(request):
    site = request.POST.get("site")
    plant = request.POST.get("plant")
    skill_category_list = []
    skill_category_id_list = []
    all_exam_books_dict = {}
    question_bank_list = QuestionBank.objects.filter(
        Site=site, Plant=plant, IsExercised='True'
    ).values()

    if question_bank_list.exists() is False:
        return JsonResponse({"question_bank_list": [], "all_exam_books_dict": {}})

    for question_bank in question_bank_list:
        skill_category_list.append(question_bank['Name'])
        skill_category_id_list.append(question_bank['QuestionBankID'])
        all_exam_books_dict[question_bank['Name']] = []

    exam_book_obj_list = Book.objects.filter(
        Site=site, Plant=plant, TrainingBankName__in=skill_category_list
    )
    if exam_book_obj_list.exists() is False:
        return JsonResponse(
            {
                "question_bank_list": list(question_bank_list),
                "all_exam_books_dict": all_exam_books_dict,
            }
        )

    for exam_book_obj in exam_book_obj_list:
        all_exam_books_dict[exam_book_obj.TrainingBankName].append(
            model_to_dict(exam_book_obj)
        )

    return JsonResponse(
        {
            "question_bank_list": list(question_bank_list),
            "all_exam_books_dict": all_exam_books_dict,
        }
    )


# 训练，判断是否能生成试卷，可以则返回试题


@csrf_exempt
def make_practice_paper(request):
    # 获取试卷属性
    q_bankID = request.POST.get('Q_bankID', '')
    book_id = request.POST.get('Q_BookID', '')
    question_num = request.POST.get('question_num', 0)
    option_ans_categorys = request.POST.getlist('NG_category_list', [])
    option_ans_min_num = request.POST.get('optionAns_min_num', 0)
    new_q_beforeDays = request.POST.get('newQ_before_days', 0)
    new_q_start_date = timezone.now() - timezone.timedelta(days=int(new_q_beforeDays))
    new_q_percent = request.POST.get('newQ_percent', 0)
    # 判断试卷是否能生成
    judge = get_exam_questions(
        int(q_bankID),
        int(book_id),
        int(question_num),
        int(option_ans_min_num),
        new_q_start_date,
        int(new_q_percent),
        option_ans_categorys,
    )
    # print('试卷是否能生成------------------\nresult_code\n',judge['result_code'],'\nresult_text\n',judge['result_text'],'\nbook_question_num\n',judge['book_question_num'],'\nexam_questions\n',judge['exam_questions'],'\n返回考题的类型:\n',type(judge['exam_questions']))
    if judge['result_code'] != 0:
        return JsonResponse(
            {"result_code": judge['result_code'], "result_text": judge['result_text']}
        )
    # 生成题目和选项
    qID_list = []
    option_dict = {}
    for question in judge['exam_questions']:
        qID_list.append(question['ChoiceQuestionID'])
        option_dict[question['ChoiceQuestionID']] = []
        # 2.通过题目ID查找该题的所有选项
        option_list = QuestionOption.objects.filter(ChoiceQuestionID__in=qID_list)
    # print('\noption_list',option_list)
    if option_list.exists() is False:
        return JsonResponse(
            {
                "result_code": judge['result_code'],
                "result_text": judge['result_text'],
                "practice_questions": judge['exam_questions'],
                "option_dict": {},
            }
        )
    for option_obj in option_list:
        option_dict[option_obj.ChoiceQuestionID].append(model_to_dict(option_obj))
    # print(option_dict)
    return JsonResponse(
        {
            "result_code": judge['result_code'],
            "result_text": judge['result_text'],
            "practice_questions": judge['exam_questions'],
            "option_dict": option_dict,
        }
    )


@csrf_exempt
def get_all_exam_books(request):
    skill_category_list = []
    all_exam_books_dict = {}

    staff_id = request.POST.get('Staff_ID', '')
    site = request.POST.get("site")
    plant = request.POST.get("plant")
    role = get_user_role(staff_id)
    if role is None:
        return JsonResponse({"question_bank_list": [], "all_exam_books_dict": {}})
    elif role == "Admin":
        question_bank_list = QuestionBank.objects.filter(
            Site=site, Plant=plant, IsExercised='False'
        ).values()
    else:
        question_bank_list = QuestionBank.objects.filter(
            Site=site,
            Plant=plant,
            ExamPaperPublisher__contains=staff_id,
            IsExercised='False',
        ).values()

    if question_bank_list.exists() is False:
        return JsonResponse({"question_bank_list": [], "all_exam_books_dict": {}})

    for question_bank in question_bank_list:
        skill_category_list.append(question_bank['Name'])
        all_exam_books_dict[question_bank['Name']] = []

    exam_book_obj_list = Book.objects.filter(
        Site=site, Plant=plant, TrainingBankName__in=skill_category_list
    )
    if exam_book_obj_list.exists() is False:
        return JsonResponse(
            {
                "question_bank_list": list(question_bank_list),
                "all_exam_books_dict": all_exam_books_dict,
            }
        )

    for exam_book_obj in exam_book_obj_list:
        all_exam_books_dict[exam_book_obj.TrainingBankName].append(
            model_to_dict(exam_book_obj)
        )

    # print(all_exam_books_dict)
    return JsonResponse(
        {
            "question_bank_list": list(question_bank_list),
            "all_exam_books_dict": all_exam_books_dict,
        }
    )


# 考核，查询题库所有考核技能类型及教材
@csrf_exempt
def get_exam_skill_type(request):
    exam_books = []
    skill_type = []
    exam_books_dict = {}

    staff_id = request.POST.get('Staff_ID', '')
    role = get_user_role(staff_id)
    if role is None:
        return JsonResponse({"skill_type": list(skill_type)})
    elif role == "Admin":
        skill_type = QuestionBank.objects.filter(IsExercised='False').values()
    else:
        skill_type = QuestionBank.objects.filter(
            ExamPaperPublisher__contains=staff_id, IsExercised='False'
        ).values()

    return JsonResponse({"skill_type": list(skill_type)})


# 考核，查询题库所有考核技能类型及教材


@csrf_exempt
def get_training_book(request):
    training_bank_id = request.POST.get('training_bank_id', '')
    training_bank_name = request.POST.get('training_bank_name', '')

    exam_books = Book.objects.filter(
        TrainingBankID=training_bank_id, TrainingBankName=training_bank_name
    ).values()
    if exam_books.exists():
        return JsonResponse({"exam_books": list(exam_books)})
    else:
        return JsonResponse({"exam_books": []})


# 考核，检查考生site和厂别信息是否符合
@csrf_exempt
def judge_exam_site_plant(request):
    site = request.POST.get('site')
    plant = request.POST.get('plant')
    examinees = request.POST.get('examinee')
    all_right = []
    not_same_site_plant = []
    no_register = []
    # 字符串转字典
    examinees = eval(examinees)
    for examinee in examinees:
        existExaminee = User.objects.filter(StaffID=examinee['工号']).first()
        # 考生已注册
        if existExaminee is not None:
            thisExaminee = model_to_dict(existExaminee)
            examinee['Site'] = thisExaminee['Site']
            examinee['厂别'] = thisExaminee['PlantName']
            # 已注册但不属于该Site和厂别的名单
            if (thisExaminee['Site'] != site) or (thisExaminee['PlantName'] != plant):
                not_same_site_plant.append(examinee)
            # 已注册且Site和厂别信息相符的名单
            else:
                all_right.append(examinee)
        # 考生未注册
        else:
            # 没有注册的名单
            examinee['Site'] = ''
            examinee['厂别'] = ''
            no_register.append(examinee)
    return JsonResponse(
        {
            'not_same_site_plant': not_same_site_plant,
            'all_right': all_right,
            'no_register': no_register,
        }
    )


# 考核，自动注册
@csrf_exempt
def auto_register(request):
    exam_paperID = request.POST.get('exam_paperID')
    examinees = request.POST.get('examinee')
    # bookID = request.POST.get('bookID')
    site = request.POST.get('site')
    plant = request.POST.get('plant')
    lang = request.POST.get('lang')

    # 字符串转字典
    examinees = eval(examinees)
    for examinee in examinees:
        # 若该试卷存在
        examp_paper_obj = ExamPaper.objects.filter(ExamPaperID=exam_paperID).first()
        if examp_paper_obj is not None:
            # 向Examinee存储该试卷的考核名单
            newExaminee = Examinee.objects.create(
                ExamPaperID=exam_paperID,
                StaffID=examinee['工号'],
                StaffName=examinee['姓名'],
                ProductLine=examinee['线别'],
                CreateDate=timezone.now(),
            )
            newExaminee.save()

        existExaminee = User.objects.filter(StaffID=examinee['工号'])
        # 考生未注册
        if len(existExaminee) <= 0:
            newuser = User(
                Site=site,
                PlantName=plant,
                StaffID=examinee['工号'],
                Password=examinee['工号'],
                StaffName=examinee['姓名'],
                Language=lang,
                IsActivated='1',
                Role='User',
            )
            newuser.save()
    return JsonResponse({'res': '成功'})


def get_exam_questions(
    question_bank_id,
    book_id,
    exam_question_num,
    optionAns_min_num,
    newQ_start_date,
    newQ_percent,
    exam_option_list,
):
    """
    return:
    result:
    0: it is successful;
    1: The book questions are not enough to exam.
    2: The new part exam question number is bigger than exam question number;
    3: There isn't NG type to exam.
    4: The %s NG type %s can't all show up at least %s times.
    """

    exam_question_id_list = []

    # 查看该教材试题是否足够
    choice_question_obj_list = ChoiceQuestion.objects.filter(
        QuestionBankID=question_bank_id, BookID=book_id
    )
    book_question_num = choice_question_obj_list.count()
    if book_question_num < exam_question_num:
        result_text = (
            'The book question number %s is less than exam question number %s.'
            % (book_question_num, exam_question_num)
        )
        return {
            'result_code': 1,
            'result_text': result_text,
            'book_question_num': book_question_num,
            'exam_questions': None,
        }

    new_choice_question_obj_list = choice_question_obj_list.filter(
        UpdateDate__gte=newQ_start_date
    )
    newQ_num = new_choice_question_obj_list.count()
    oldQ_num = book_question_num - newQ_num
    # print("newQ_num:%s, oldQ_num:%s" % (newQ_num, oldQ_num))

    newQ_exam_num = int(newQ_percent / 100 * newQ_num)
    if newQ_exam_num <= exam_question_num:
        # 确认旧题题数
        if exam_question_num - newQ_exam_num > oldQ_num:
            # 旧题不够
            oldQ_exam_num = oldQ_num
            newQ_exam_num = exam_question_num - oldQ_exam_num
        else:
            oldQ_exam_num = exam_question_num - newQ_exam_num
    else:
        newQ_exam_num = exam_question_num
        oldQ_exam_num = 0
        # result_text = 'The new part exam question number %s is bigger than exam question number %s.' % (newQ_exam_num, exam_question_num)
        # return {'result_code': 2, 'result_text': result_text, 'book_question_num': book_question_num, 'exam_questions': None}

    # print("exam_question_num:%s, newQ_exam_num:%s, oldQ_exam_num:%s" % (exam_question_num, newQ_exam_num, oldQ_exam_num))

    # 捞取所有需要考核的NG类型
    # option_list = Option.objects.all().values_list("Name", flat=True)
    # exam_option_list = list(option_list)

    if not exam_option_list:
        result_text = "There isn't NG type to exam."
        return {
            'result_code': 3,
            'result_text': result_text,
            'book_question_num': book_question_num,
            'exam_questions': None,
        }

    unfinished_option_list = list(set(exam_option_list))
    print(exam_option_list)
    print(unfinished_option_list, '123')
    question_option_selected_num_dict = {}
    for option in unfinished_option_list:
        question_option_selected_num_dict[option] = 0

    new_question_option_obj_list = QuestionOption.objects.filter(
        ChoiceQuestionID__in=new_choice_question_obj_list, IsAnswer=True
    )
    new_exam_question_id_list = []
    question_option_selected_id_list = []
    while len(new_exam_question_id_list) < newQ_exam_num:
        # 在没有完成获取的NG类型试题选项中，去除选过的试题选项，随机一个试题选项
        tmp_question_option_obj_list = (
            new_question_option_obj_list.filter(Option__in=unfinished_option_list)
            .exclude(QuestionOptionID__in=question_option_selected_id_list)
            .order_by('?')[:1]
        )
        if tmp_question_option_obj_list.exists():
            question_id = tmp_question_option_obj_list[0].ChoiceQuestionID

            if question_id not in new_exam_question_id_list:
                new_exam_question_id_list.append(question_id)

            tmp_question_option_obj_list = new_question_option_obj_list.filter(
                ChoiceQuestionID=question_id
            )
            for question_option_obj in tmp_question_option_obj_list:
                question_option_id = question_option_obj.QuestionOptionID
                question_option_selected_id_list.append(question_option_id)

                option = question_option_obj.Option
                if option in unfinished_option_list:
                    question_option_selected_num_dict[option] += 1
                    # 去除完成NG类型至少出现N次的NG类型
                    if question_option_selected_num_dict[option] >= optionAns_min_num:
                        unfinished_option_list.remove(option)

            # 所有NG类型至少出现N次的条件满足后，直接结束循环。
            if not unfinished_option_list:
                break
        else:
            # 已经没有适合的试题选项可以捞取了，直接结束循环。
            break

    num = newQ_exam_num - len(new_exam_question_id_list)
    if num > 0:
        # 随机要求新题部分的百分比试题数的不足数量
        tmp_question_id_list = (
            new_choice_question_obj_list.exclude(
                ChoiceQuestionID__in=new_exam_question_id_list
            )
            .values_list('ChoiceQuestionID', flat=True)
            .order_by('?')[:num]
        )
        new_exam_question_id_list.extend(list(tmp_question_id_list))

    # print(new_exam_question_id_list)
    # print(question_option_selected_num_dict)
    # print(question_option_selected_id_list)

    old_choice_question_obj_list = choice_question_obj_list.filter(
        UpdateDate__lt=newQ_start_date
    )
    old_question_option_obj_list = QuestionOption.objects.filter(
        ChoiceQuestionID__in=old_choice_question_obj_list, IsAnswer=True
    )
    question_option_selected_id_list = []
    old_exam_question_id_list = []

    while len(old_exam_question_id_list) < oldQ_exam_num:
        # 在没有完成获取的NG类型试题选项中，去除选过的试题选项，随机一个试题选项
        tmp_question_option_obj_list = (
            old_question_option_obj_list.filter(Option__in=unfinished_option_list)
            .exclude(QuestionOptionID__in=question_option_selected_id_list)
            .order_by('?')[:1]
        )
        if tmp_question_option_obj_list.exists():
            question_id = tmp_question_option_obj_list[0].ChoiceQuestionID

            if question_id not in old_exam_question_id_list:
                old_exam_question_id_list.append(question_id)

            tmp_question_option_obj_list = old_question_option_obj_list.filter(
                ChoiceQuestionID=question_id
            )
            for question_option_obj in tmp_question_option_obj_list:
                question_option_id = question_option_obj.QuestionOptionID
                question_option_selected_id_list.append(question_option_id)

                option = question_option_obj.Option
                if option in unfinished_option_list:
                    question_option_selected_num_dict[option] += 1
                    # 去除完成NG类型至少出现N次的NG类型
                    if question_option_selected_num_dict[option] >= optionAns_min_num:
                        unfinished_option_list.remove(option)

            # 所有NG类型至少出现N次的条件满足后，直接结束循环。
            if not unfinished_option_list:
                break
        else:
            # 已经没有适合的试题选项可以捞取了，直接结束循环。
            break

    # 在旧题部分中随机获取剩余的考题
    num = oldQ_exam_num - len(old_exam_question_id_list)
    if num > 0:
        tmp_question_id_list = (
            old_choice_question_obj_list.exclude(
                ChoiceQuestionID__in=old_exam_question_id_list
            )
            .values_list('ChoiceQuestionID', flat=True)
            .order_by('?')[:num]
        )
        old_exam_question_id_list.extend(list(tmp_question_id_list))

    # print(old_exam_question_id_list)
    # print(question_option_selected_num_dict)

    print(unfinished_option_list)
    if unfinished_option_list:
        result_text = "The %s NG type %s can't all show up at least %s times." % (
            len(unfinished_option_list),
            ",".join(unfinished_option_list),
            optionAns_min_num,
        )
        return {
            'result_code': 4,
            'result_text': result_text,
            'book_question_num': book_question_num,
            'exam_questions': None,
        }

    exam_question_id_list = new_exam_question_id_list + old_exam_question_id_list
    # print(exam_question_id_list)

    exam_question_obj_list = choice_question_obj_list.filter(
        ChoiceQuestionID__in=exam_question_id_list
    ).values()

    result_text = 'it is successful.'
    return {
        'result_code': 0,
        'result_text': result_text,
        'book_question_num': book_question_num,
        'exam_questions': list(exam_question_obj_list),
    }


# 考核，发布试卷，存储新试卷属性
@csrf_exempt
def to_issue_paper(request):
    # 试卷属性注意：试卷名不能重复，符合条件的题目数量要够等……及必须确保该试卷可以合成
    # 发布试卷上传的名单去掉重复的工号
    creator_site = request.POST.get('creator_site', '')
    creator_plant = request.POST.get('creator_plant', '')
    paper_name = request.POST.get('paper_name', '')
    Q_bank_id = request.POST.get('Q_bankID', '')
    Q_bank_name = request.POST.get('Q_bankName', '')
    book_id = request.POST.get('Q_BookID', '')
    IsExercised = request.POST.get('IsExercised', 'True')
    exam_question_num = request.POST.get('question_num', 0)
    pass_question_num = request.POST.get('pass_question_num', 0)
    Q_limit_time = request.POST.get('Q_limit_time', 0)
    exam_time = int(exam_question_num) * int(Q_limit_time)
    optionAns_min_num = request.POST.get('optionAns_min_num', 0)
    newQ_before_days = request.POST.get('newQ_before_days', 0)
    newQ_percent = request.POST.get('newQ_percent', 0)
    deadline_utcstamp_millisec = request.POST.get('deadline_utcstamp_millisec', 0)
    exam_max_count = request.POST.get('exam_max_count', 0)
    creator_id = request.POST.get('creator_id', '')
    creator_name = request.POST.get('creator_name', '')
    NG_category_list = request.POST.getlist('NG_category_list', [])

    # print(NG_category_list)

    exam_paper_obj_list = ExamPaper.objects.filter(
        Site=creator_site, Plant=creator_plant, Name=paper_name
    )
    if exam_paper_obj_list.exists():
        result_text = 'The Exam Paper "%s" has existed!' % (paper_name)
        context = {
            'result_code': 5,
            'result_text': result_text,
            'book_question_num': 0,
            'exam_questions': None,
        }
        return JsonResponse(context)

    exam_deadline = datetime.utcfromtimestamp(
        int(deadline_utcstamp_millisec) / 1000
    ).replace(tzinfo=timezone.utc)

    # print('试卷名称:',paper_name,'题库:',Q_bank_name,'题库id:',Q_bank_id, '教材ID:', book_id, '用于训练:', IsExercised, '试卷题目数量:',exam_question_num,'考核通过题目数量:',
    # pass_question_num,'每道题限制时间:',Q_limit_time,'考试时间:',exam_time,'新题的时间段(近数天):',
    # newQ_before_days,'新题需占百分比:',newQ_percent,'不良类型至少出现次数:',optionAns_min_num,'考试截止日期:',exam_deadline, '试卷最大考核次数:',exam_max_count,'发布者ID:',creator_id, '发布者名字:',creator_name)

    if IsExercised == 'False':
        IsExercised = False
    else:
        IsExercised = True

    newQ_start_date = timezone.now() - timezone.timedelta(days=int(newQ_before_days))
    # newQ_start_date = datetime.strptime('2021-1-5 08:00:00 +0800', '%Y-%m-%d %H:%M:%S %z').astimezone(tz=timezone.utc)
    # print(newQ_start_date)

    ret = get_exam_questions(
        Q_bank_id,
        book_id,
        int(exam_question_num),
        int(optionAns_min_num),
        newQ_start_date,
        int(newQ_percent),
        NG_category_list,
    )
    # print(ret)

    if ret['result_code'] == 0:
        exam_paper_obj = ExamPaper.objects.create(
            Site=creator_site,
            Plant=creator_plant,
            Name=paper_name,
            QuestionBankID=Q_bank_id,
            QuestionBankName=Q_bank_name,
            IsExercised=IsExercised,
            BookID=book_id,
            Deadline=exam_deadline,
            ExamMaxCount=exam_max_count,
            QuestionNumber=exam_question_num,
            PassQuestionNumber=pass_question_num,
            QuestionLimitTime=Q_limit_time,
            ExamTime=exam_time,
            OptionAnsCategorys=','.join(NG_category_list),
            OptionAnsMinNumber=optionAns_min_num,
            NewQuestionBeforeDays=newQ_before_days,
            NewQuestionPercent=newQ_percent,
            CreatorID=creator_id,
            CreatorName=creator_name,
            CreateDate=timezone.now(),
        )
        ret['exam_paper_id'] = exam_paper_obj.ExamPaperID

    # print('自增id:',newpaper.ExamPaperID)
    # return JsonResponse({"exam_paperID":newpaper.ExamPaperID,'exam_deadLine':exam_deadline,'exam_max_count':exam_max_count,'res':'成功'})
    return JsonResponse(ret)


# 考核，查询具体发布者的所有考核


@csrf_exempt
def get_all_issue(req):
    staffID = req.POST['staffID']
    all_issue_list = ExamPaper.objects.filter(CreatorID=staffID).values()
    if all_issue_list.exists() is False:
        return JsonResponse({"all_issue": []})

    exam_paperID_list = []
    examinee_dict = {}
    for obj in all_issue_list:
        obj['bookName'] = '该教材已被删除'

        book_list = Book.objects.filter(BookID=obj['BookID']).values('Name')
        if book_list.exists() is True:
            obj['bookName'] = book_list[0]['Name']

        examinee_count = Examinee.objects.filter(ExamPaperID=obj['ExamPaperID']).count()
        obj['examineeCount'] = examinee_count

        examinee_pass = Examinee.objects.filter(ExamPaperID=obj['ExamPaperID']).filter(
            IsPassed=False
        )
        if len(examinee_pass) > 0:
            obj['Examineepass'] = 'True'
        else:
            obj['Examineepass'] = 'False'

    return JsonResponse({"all_issue": list(all_issue_list)})


# 考核，查询具体考生所有考核


@csrf_exempt
def get_all_exam(req):
    staffID = req.POST['staffID']
    # 考核名单中找该staffID的记录
    examinee_list = Examinee.objects.filter(StaffID=staffID).values()
    exam_paper_list = ExamPaper.objects.all().values()
    all_examed_paper_list = ExamedPaper.objects.filter(ExamineeStaffID=staffID).values(
        'ExamPaperID', 'ExamedCount'
    )
    # 如果Examinee中有该staffID的记录
    if examinee_list:
        # Examinee中该staffID的所有记录转DataFrame
        examineeDF = pd.DataFrame(examinee_list)
        # 所有的试卷转DataFrame
        paperDF = pd.DataFrame(exam_paper_list)
        # 合并Examinee和Paper表
        examinee_paperDF = pd.merge(examineeDF, paperDF, how='inner', on='ExamPaperID')
        examinee_paper_list = examinee_paperDF.to_dict('records')
        for obj in examinee_paper_list:
            obj['bookName'] = '该教材已被删除'
            book_list = Book.objects.filter(BookID=obj['BookID']).values('Name')
            if book_list.exists() is True:
                obj['bookName'] = book_list[0]['Name']
        examinee_paper_bookDF = pd.DataFrame(examinee_paper_list)
        all_examed_paperDF = pd.DataFrame(all_examed_paper_list).rename(
            columns={'ExamedCount': 'ExamedCount_inExamed'}
        )
        # print('examinee_paperDF----------------------\n',examinee_paperDF)
        # print('examinee_paper_bookDF----------------------\n',examinee_paper_bookDF)
        # print('all_examed_paperDF----------------------\n',all_examed_paperDF)
        # 该用户有过任一考核记录，直接和ExamedPaper表连接
        if all_examed_paper_list:
            all_examinee_paper_book_examed_paperDF = pd.merge(
                examinee_paper_bookDF, all_examed_paperDF, how='outer', on='ExamPaperID'
            ).fillna(0)
            examinee_paper_book_examed_paperDF = all_examinee_paper_book_examed_paperDF[
                all_examinee_paper_book_examed_paperDF['ExamedCount_inExamed']
                == all_examinee_paper_book_examed_paperDF['ExamedCount']
            ]
            # print('1',all_examinee_paper_book_examed_paperDF.loc[1,:])
            # print('1该用户有过任一考核记录\n',examinee_paper_book_examed_paperDF)
        else:
            examinee_paper_book_examed_paperDF = examinee_paper_bookDF
            # print('2该用户无任何考核记录\n',examinee_paper_book_examed_paperDF)
        all_exam_paper = examinee_paper_book_examed_paperDF.to_dict('records')
        return JsonResponse({"isExaminee": "true", "all_exam_paper": all_exam_paper})
    else:
        # Examinee中没有该staffID的记录
        return JsonResponse({"isExaminee": "false", "all_exam_paper": []})


# 考核，生成试题


@csrf_exempt
def make_paper(req):
    # 获取试卷属性
    examineeID = req.POST['examineeID']
    examinee = Examinee.objects.filter(ExamineeID=examineeID)
    examinee_list = examinee.values()
    exam_paperID = examinee_list[0]['ExamPaperID']
    exam_paper_list = ExamPaper.objects.filter(ExamPaperID=exam_paperID).values()
    # print('\n该考卷的所有属性--------------------------',exam_paper_list)
    book_site = exam_paper_list[0]['Site']
    book_plant = exam_paper_list[0]['Plant']
    book_id = exam_paper_list[0]['BookID']
    exam_name = exam_paper_list[0]['Name']
    q_bankID = exam_paper_list[0]['QuestionBankID']
    q_bankName = exam_paper_list[0]['QuestionBankName']
    question_num = exam_paper_list[0]['QuestionNumber']
    pass_q_num = exam_paper_list[0]['PassQuestionNumber']
    q_limit_time = exam_paper_list[0]['QuestionLimitTime']
    exam_time = exam_paper_list[0]['ExamTime']
    option_ans_categorys = exam_paper_list[0]['OptionAnsCategorys'].split(',')
    option_ans_min_num = exam_paper_list[0]['OptionAnsMinNumber']
    new_q_beforeDays = exam_paper_list[0]['NewQuestionBeforeDays']
    new_q_start_date = timezone.now() - timezone.timedelta(days=int(new_q_beforeDays))
    new_q_percent = exam_paper_list[0]['NewQuestionPercent']
    deadline = exam_paper_list[0]['Deadline']
    exam_max_count = exam_paper_list[0]['ExamMaxCount']
    examinee_staffID = examinee_list[0]['StaffID']
    examinee_staff_name = examinee_list[0]['StaffName']
    examed_count = examinee_list[0]['ExamedCount']
    create_date = timezone.now()
    # 判断试卷是否能生成
    judge = get_exam_questions(
        q_bankID,
        book_id,
        question_num,
        option_ans_min_num,
        new_q_start_date,
        new_q_percent,
        option_ans_categorys,
    )
    # print('试卷是否能生成------------------\nresult_code\n',judge['result_code'],'\nresult_text\n',judge['result_text'],'\nbook_question_num\n',judge['book_question_num'],'\nexam_questions\n',judge['exam_questions'],'\n返回考题的类型:\n',type(judge['exam_questions']))
    if judge['result_code'] != 0:
        return JsonResponse(
            {"result_code": judge['result_code'], "result_text": judge['result_text']}
        )
    # 更新Examinee的ExamedCount
    examinee_item = examinee[0]
    examinee_item.ExamedCount += 1
    examinee_item.save()
    # 生成一条ExamedPaper记录
    newExamedPaper = ExamedPaper(
        Site=book_site,
        Plant=book_plant,
        ExamPaperID=exam_paperID,
        Name=exam_name,
        QuestionBankID=q_bankID,
        QuestionBankName=q_bankName,
        QuestionNumber=question_num,
        PassQuestionNumber=pass_q_num,
        QuestionLimitTime=q_limit_time,
        ExamTime=exam_time,
        OptionAnsMinNumber=option_ans_min_num,
        NewQuestionBeforeDays=new_q_beforeDays,
        NewQuestionPercent=new_q_percent,
        Deadline=deadline,
        ExamMaxCount=exam_max_count,
        ExamineeStaffID=examinee_staffID,
        ExamineeStaffName=examinee_staff_name,
        ExamedCount=(examed_count + 1),
        PassedQuestionNumber=0,
        CreateDate=create_date,
    )
    newExamedPaper.save()
    examed_paperID = newExamedPaper.ExamedPaperID
    # 生成题目
    for question in judge['exam_questions']:
        # 1.生成对应ExamedPaper记录的ExamedChoiceQuestion题目
        newExamedChoiceQuestion = ExamedChoiceQuestion(
            ExamedPaperID=examed_paperID,
            ChoiceQuestionID=question['ChoiceQuestionID'],
            Subject=question['Subject'],
            Type=question['Type'],
            Content=question['Content'],
        )
        newExamedChoiceQuestion.save()
        examed_questionID = newExamedChoiceQuestion.ExamedChoiceQuestionID
        # 2.通过题目ID查找该题的所有选项
        options = QuestionOption.objects.filter(
            ChoiceQuestionID=question['ChoiceQuestionID']
        )
        # 3.生成对应ExamedPaper记录的ExamedQuestionOption选项
        for option in options:
            newExamedQuestionOption = ExamedQuestionOption(
                ExamedChoiceQuestionID=examed_questionID,
                Option=option.Option,
                IsAnswer=option.IsAnswer,
            )
            newExamedQuestionOption.save()
    # print('\n新生成的examed_paperID自增ID',examed_paperID)
    return JsonResponse(
        {
            "result_code": judge['result_code'],
            "result_text": judge['result_text'],
            "exam_questions": judge['exam_questions'],
            "examed_paperID": examed_paperID,
        }
    )


# 考核，通过examed_paperID查找具体ExamedPaper的试题


@csrf_exempt
def get_paper(req):
    examed_paperID = req.POST['examed_paperID']
    question_list = ExamedChoiceQuestion.objects.filter(ExamedPaperID=examed_paperID)
    option_list = []
    for question in question_list:
        options = ExamedQuestionOption.objects.filter(
            ExamedChoiceQuestionID=question.ExamedChoiceQuestionID
        ).values()
        option_list += options[:]
    return JsonResponse(
        {
            'res': 'yes',
            'question_list': list(question_list.values()),
            'option_list': list(option_list),
        }
    )


# 考核，更新作答信息


@csrf_exempt
def update_exam_record(req):
    examed_question_id = req.POST['examed_question_id']
    examed_q_optionID = req.POST['examed_q_optionID']
    examed_question_list = ExamedChoiceQuestion.objects.filter(
        ExamedChoiceQuestionID=examed_question_id
    )
    examed_option_list = ExamedQuestionOption.objects.filter(
        ExamedQuestionOptionID=examed_q_optionID
    )
    examed_paperID = examed_question_list.values('ExamedPaperID')[0]['ExamedPaperID']
    examed_paper_list = ExamedPaper.objects.filter(ExamedPaperID=examed_paperID)
    # 更新ExamedQuestionOption
    examed_option = examed_option_list[0]
    examed_option.IsMarked = True
    examed_option.save()
    is_answer = examed_option_list.values('IsAnswer')[0]['IsAnswer']
    # 选择的选项是正确答案，则更新ExamedPaper
    if is_answer == True:
        examed_paper = examed_paper_list[0]
        examed_paper_values = examed_paper_list.values(
            'PassQuestionNumber', 'PassedQuestionNumber'
        )
        pass_q_num = examed_paper_values[0]['PassQuestionNumber']
        passed_q_num = examed_paper_values[0]['PassedQuestionNumber']
        # 通过的题目数量已达标
        if passed_q_num + 1 >= pass_q_num:
            examed_paper.IsPassed = True
            # 更新Examinee表的IsPassed
            examinee = Examinee.objects.filter(
                ExamPaperID=examed_paper.ExamPaperID,
                StaffID=examed_paper.ExamineeStaffID,
            )[0]
            if examinee.IsPassed == False:
                examinee.IsPassed = True
                examinee.save()
        examed_paper.PassedQuestionNumber += 1
        examed_paper.save()
    return JsonResponse({'res': 'yes'})


# 考核，查找考核记录


@csrf_exempt
def get_exam_record(request):
    exam_role = request.POST.get('role')
    admin_site = request.POST.get('admin_site')
    admin_plant = request.POST.get('admin_plant')
    exam_staff_ID = request.POST.get('staff_ID')
    result = []
    try:
        if exam_role == 'Admin':
            all_exam_record = (
                ExamedPaper.objects.filter(
                    Site=admin_site, Plant=admin_plant, IsExercised=False
                )
                .order_by('-CreateDate')
                .values()
            )
            all_exam_record = list(all_exam_record)
        else:
            all_exam_record = (
                ExamedPaper.objects.filter(ExamineeStaffID=exam_staff_ID)
                .filter(IsExercised=False)
                .order_by('-CreateDate')
                .values()
            )
            all_exam_record = list(all_exam_record)
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": all_exam_record})


# mwj now
# 考核记录详情


@csrf_exempt
def examedChoiceQuestion_page(request, examedPaperID, name, examedCount):
    context = {'examedPaperID': examedPaperID, 'name': name, 'examedCount': examedCount}
    return render(request, "examed_cq.html", context)


@csrf_exempt
def examedChoiceQuestion_page_data(request, examedPaperID, name):
    examed_cq = ExamedChoiceQuestion.objects.filter(
        ExamedPaperID=examedPaperID
    ).values()
    examed_cq = list(examed_cq)
    return JsonResponse({"examed_cq": examed_cq})


# 获得考过的题的所有选项


@csrf_exempt
def get_examed_options(req):
    req.encoding = 'utf-8'
    try:
        ExamedChoiceQuestionID = req.GET['ExamedChoiceQuestionID']
        options = ExamedQuestionOption.objects.filter(
            ExamedChoiceQuestionID=ExamedChoiceQuestionID
        )
        result = []
        # 遍历结果
        for e in options:
            # 转成字典（json类型）
            jsongstr = model_to_dict(e)
            # 塞返回结果list
            result.append(jsongstr)
        return JsonResponse({"options": result})
    except Exception as e:
        return HttpResponse(repr(e))


# 训练，查找训练记录


@csrf_exempt
def get_practice_record(request):
    practice_staff_ID = request.POST.get('staff_ID')
    try:
        all_practice_record = serializers.serialize(
            "json",
            ExamedPaper.objects.filter(ExamineeStaffID=practice_staff_ID)
            .filter(IsExercised=True)
            .order_by('-CreateDate'),
        )
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": all_practice_record})


# 注册界面
@csrf_exempt
def add_new_user(request):
    StaffID = request.POST.get('StaffID')
    StaffName = request.POST.get('StaffName')
    Password = request.POST.get('Password')
    Language = request.POST.get('Language')
    Department = request.POST.get('Department')
    newSite = request.POST.get('newSite')
    newPlant = request.POST.get('newPlant')
    repert_staffid = User.objects.filter(StaffID=StaffID)
    # print(Language)
    if len(repert_staffid) > 0:
        return JsonResponse({"result": "false"})
    else:
        # if re.fullmatch(r"[A-Z]{1}[0-9]{8}",StaffID):
        newuser = User(
            Site=newSite,
            PlantName=newPlant,
            StaffID=StaffID,
            Password=Password,
            StaffName=StaffName,
            Department=Department,
            Language=Language,
            IsActivated='1',
            Role='User',
        )
        newuser.save()
        return JsonResponse({"result": "true"})
        # else:
        #     return JsonResponse({"result":"false"})


# 用户登录验证密码


@csrf_exempt
def login_index(request):
    StaffID_login = request.POST.get('StaffID_login')
    Password_login = request.POST.get('Password_login')
    user = User.objects.filter(StaffID=StaffID_login).first()
    # logger.info('\nb info:\n{}'.format(user))
    if user:
        myinfo = User.objects.get(StaffID=StaffID_login)
        pwd = myinfo.Password
        staffName = myinfo.StaffName
        staffID = myinfo.StaffID
        staffSite = myinfo.Site
        staffPlant = myinfo.PlantName
        staffRole = myinfo.Role
        staffLang = myinfo.Language
        if Password_login == pwd:
            # request.session['StaffID'] = StaffID_login
            return JsonResponse(
                {
                    "result": "true",
                    'name': staffName,
                    'S_StaffID': staffID,
                    'S_Site': staffSite,
                    'S_Plant': staffPlant,
                    'S_Role': staffRole,
                    'S_Lang': staffLang,
                }
            )
        else:
            return JsonResponse({"result": "false"})
    else:
        return JsonResponse({"result": "false"})


# 个人信息


@csrf_exempt
def get_user(request):
    personal_user = request.POST.get('personal_ID')
    # print(personal_user)
    try:
        Personal_List = serializers.serialize(
            "json", User.objects.filter(StaffID=personal_user)
        )
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": Personal_List})


# 管理员界面操作
# 查询所有用户


@csrf_exempt
def get_all_user(request):
    admin_site = request.POST.get('admin_site')
    admin_plant = request.POST.get('admin_plant')
    try:
        UserList = serializers.serialize(
            "json",
            User.objects.filter(Site=admin_site, PlantName=admin_plant).order_by(
                'Site', 'PlantName', 'StaffID'
            ),
        )
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": UserList})


# 删除用户
@csrf_exempt
def del_user(request):
    StaffID_del = request.POST.get('StaffID')
    user_del = User.objects.get(StaffID=StaffID_del)
    user_del.delete()
    return JsonResponse({"data": "true"})


# 修改用户信息
@csrf_exempt
def modify_user(request):
    modifySite = request.POST.get('modifySite')
    modifyPlant = request.POST.get('modifyPlant')
    inputStaffID = request.POST.get('modifyStaffID')
    modifyStaffRole = request.POST.get('modifyStaffRole')
    inputPassword = request.POST.get('modifyPassword')
    inputDepartment = request.POST.get('modifyDepartment')
    inputLanguage = request.POST.get('modifyLanguage')
    # 查询出来需要需改的StaffID
    modifyUser = User.objects.get(StaffID=inputStaffID)
    if modifySite:
        modifyUser.Site = modifySite
    if modifyPlant:
        modifyUser.PlantName = modifyPlant
    if modifyStaffRole:
        modifyUser.Role = modifyStaffRole
    modifyUser.Password = inputPassword
    modifyUser.Department = inputDepartment
    modifyUser.Language = inputLanguage
    modifyUser.save()
    return JsonResponse({"data": "true"})


# 管理员界面添加用户


@csrf_exempt
def add_user(request):
    input_user_StaffID = request.POST.get('newStaffID')
    input_user_StaffName = request.POST.get('newStaffName')
    input_user_Password = request.POST.get('newPassword')
    input_user_Department = request.POST.get('newDepartment')
    input_user_Language = request.POST.get('newLanguage')
    user_Site = request.POST.get('newSite')
    user_Plant = request.POST.get('newPlant')
    r = User.objects.filter(StaffID=input_user_StaffID)
    if len(r) > 0:
        return JsonResponse({"result": "false"})
    else:
        # if re.fullmatch(r"[A-Z]{1}[0-9]{8}",StaffID):
        new1user = User(
            Site=user_Site,
            PlantName=user_Plant,
            StaffID=input_user_StaffID,
            StaffName=input_user_StaffName,
            Role='User',
            Department=input_user_Department,
            Password=input_user_Password,
            Language=input_user_Language,
            IsActivated='1',
        )
        new1user.save()
        return JsonResponse({"result": "true"})
    # else:
    #     return JsonResponse({"result":"false"})


# 管理员界面修改技能类型
# 查询所有技能类型


@csrf_exempt
def get_all_technology(request):
    technology_site = request.POST.get('admin_site')
    technology_plant = request.POST.get('admin_plant')
    try:
        TechnologyList = Technology.objects.filter(
            Site=technology_site, PlantName=technology_plant
        ).values()
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": list(TechnologyList)})


# 删除技能类型


@csrf_exempt
def del_technology(request):
    TechnologyID = request.POST.get('TechnologyID')
    # Technology_del = request.POST.get('Name')
    # Technology_site = request.POST.get('Site')
    # Technology_plant = request.POST.get('Plant')
    technology_name_del = Technology.objects.get(TechnologyID=TechnologyID)
    technology_name_del.delete()
    return JsonResponse({"data": "true"})


# 修改技能类型名称


@csrf_exempt
def modify_technology(request):
    # 先通过inputtechnology查到所要修改的用户，修改该用户的其他属性
    inputtechnology = request.POST.get('Name')
    # dict ['changename'] = inputtechnology
    # logger.info('\ninputtechnology:\n{}'.format(inputtechnology))
    modifyTechnology = Technology.objects.get(Name=inputtechnology)
    # 获得修改后的数据
    modifyTechnology.Name = inputtechnology
    modifyTechnology.save()
    return JsonResponse({"data": "true"})


# 管理员界面添加技能类型


@csrf_exempt
def add_technology(request):
    input_technology_Name = request.POST.get('newtechnology')
    technology_site = request.POST.get('newsite')
    technology_plant = request.POST.get('newplant')
    t = Technology.objects.filter(
        Site=technology_site, PlantName=technology_plant, Name=input_technology_Name
    )
    if len(t) > 0:
        return JsonResponse({"result": "false"})
    else:
        new_technology_ID = Technology(
            Site=technology_site,
            PlantName=technology_plant,
            Name=input_technology_Name,
            Description='',
        )
        new_technology_ID.save()
        return JsonResponse({"result": "true"})


# 管理员界面修改厂别及语言
# 查询所有厂别语言


@csrf_exempt
def get_all_plant(req):
    try:
        PlantList = serializers.serialize("json", Plant.objects.all())
    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": PlantList})


# 删除技能类型


# @csrf_exempt
# def del_plant(request):
#     plant_del = request.POST.get('ShortName')
#     plant_name_del = Plant.objects.get(ShortName=plant_del)
#     plant_name_del.delete()
#     return JsonResponse({"data": "true"})

# 修改技能类型名称


# @csrf_exempt
# def modify_plant(request):
#     inputShortName = request.POST.get('modifyShortName')
#     inputshortlanguage = request.POST.get('modifyplantLanguage')
#     modifyPlant = Plant.objects.get(ShortName=inputShortName)
#     modifyPlant.Language = inputshortlanguage
#     modifyPlant.save()
#     return JsonResponse({"data": "true"})

# 管理员界面添加技能类型


# @csrf_exempt
# def add_plant(request):
#     input_newShortName = request.POST.get('newShortName')
#     input_newplantlanguage = request.POST.get('newplantLanguage')
#     p = Plant.objects.filter(ShortName=input_newShortName)
#     if len(p) > 0:
#         return JsonResponse({"result": "false"})
#     else:
#         new_Plant_ID = Plant(ShortName=input_newShortName,
#                              Language=input_newplantlanguage)
#         new_Plant_ID.save()
#         return JsonResponse({"result": "true"})

# 管理员界面修改题库与角色
# 查询所有题库与角色


@csrf_exempt
def get_all_questionbank(request):
    # try:
    #     QuestionBankList= serializers.serialize("json",QuestionBank.objects.all())
    #     print(QuestionBankList)
    # except Exception as e:
    #     return HttpResponse(repr(e))
    # return JsonResponse({"data":QuestionBankList})
    try:
        result = []
        # staff为空会捞出全部，否则捞出具体Staff_ID相关的数据。
        staff_id = request.GET.get('Staff_ID', '')
        Admin_Site = request.GET.get('Site', '')
        Admin_Plant = request.GET.get('Plant', '')
        if staff_id == "":
            return JsonResponse({"data": result})

        user_obj = User.objects.filter(StaffID=staff_id).first()
        if user_obj is None:
            return JsonResponse({"data": result})

        if user_obj.Role == "Admin":
            QuestionBankList = QuestionBank.objects.filter(
                Site=Admin_Site, Plant=Admin_Plant
            )
        else:
            QuestionBankList = QuestionBank.objects.filter(
                Q(ExamQuestionMaker__contains=staff_id)
                | Q(ExamPaperPublisher__contains=staff_id)
            )

        # 遍历结果
        for e in QuestionBankList:
            # 转成字典（json类型）
            jsongstr = model_to_dict(e)
            # 塞返回结果list
            result.append(jsongstr)

    except Exception as e:
        return HttpResponse(repr(e))
    return JsonResponse({"data": result})


# 删除题库与角色


@csrf_exempt
def del_questionbank(request):
    get_choicequestions = request.POST.get('QuestionBankID')
    # questionbank_del_name = request.POST.get('Name')
    # questionbank_del_isexc = request.POST.get('IsExercised')
    # print(get_choicequestions,questionbank_del_name,questionbank_del_isexc)

    dq = ChoiceQuestion.objects.filter(QuestionBankID=get_choicequestions)
    # print(len(dq))
    if len(dq) > 0:
        return JsonResponse({"result": "false"})
    else:
        questionbank_name_del = QuestionBank.objects.filter(
            QuestionBankID=get_choicequestions
        )
        # print('questionbank_name_del=',questionbank_name_del)
        questionbank_name_del.delete()
        return JsonResponse({"result": "true"})


# 修改题库与角色
@csrf_exempt
def modify_questionbank(request):
    QuestionBankID = request.POST.get('QuestionBankID')
    # input_questionname = request.POST.get('modifyquestionname')
    input_isexercised = request.POST.get('modifyisexercised')
    input_examquestionmaker = request.POST.get('modifyexamquestionmaker')
    input_exampaperpublisher = request.POST.get('modifyexampaperpublisher')
    modifyQuestionBank = QuestionBank.objects.get(QuestionBankID=QuestionBankID)
    modifyQuestionBank.ExamQuestionMaker = input_examquestionmaker
    modifyQuestionBank.ExamPaperPublisher = input_exampaperpublisher
    modifyQuestionBank.save()
    return JsonResponse({"data": "true"})


# 管理员界面添加题库与角色
@csrf_exempt
def add_questionbank(request):
    site = request.POST.get('newquestionsite')
    plant = request.POST.get('newquestionplant')
    inputquestionname = request.POST.get('newquestionname')
    inputisexercised = request.POST.get('newisexercised')
    inputexamquestionmaker = request.POST.get('newexamquestionmaker')
    inputexampaperpublisher = request.POST.get('newexampaperpublisher')
    q = QuestionBank.objects.filter(
        Site=site, Plant=plant, Name=inputquestionname, IsExercised=inputisexercised
    )  # 两层筛选
    if len(q) > 0:
        return JsonResponse({"result": "false"})
    else:
        new_QuestionBank_ID = QuestionBank(
            Site=site,
            Plant=plant,
            Name=inputquestionname,
            IsExercised=inputisexercised,
            ExamQuestionMaker=inputexamquestionmaker,
            ExamPaperPublisher=inputexampaperpublisher,
        )
        new_QuestionBank_ID.save()
        return JsonResponse({"result": "true"})


# 管理员界面修改教材库与角色
# 查询所有教材库与角色
@csrf_exempt
def get_all_trainingbank(request):
    try:
        result = []
        # staff为空会捞出全部，否则捞出具体Staff_ID相关的数据。
        staff_id = request.GET.get('Staff_ID', '')
        admin_site = request.GET.get('site')
        admin_plant = request.GET.get('plant')
        if staff_id == "":
            return JsonResponse({"data": result})

        user_obj = User.objects.filter(StaffID=staff_id).first()
        if user_obj is None:
            return JsonResponse({"data": result})

        if user_obj.Role == "Admin":
            TrainingBankList = TrainingBank.objects.filter(
                Site=admin_site, Plant=admin_plant
            )
        else:
            TrainingBankList = TrainingBank.objects.filter(
                Q(TrainingMaterialMaker__contains=staff_id)
                | Q(Trainer__contains=staff_id)
            )

        # 遍历结果
        for e in TrainingBankList:
            # 转成字典（json类型）
            jsongstr = model_to_dict(e)
            # 塞返回结果list
            result.append(jsongstr)
    except Exception as e:
        return HttpResponse(repr(e))

    # print(result)
    return JsonResponse({"data": result})


# 删除教材库与角色
@csrf_exempt
def del_trainingbank(request):
    TrainingBankID = request.POST.get('TrainingBankID')
    # trainingbank_del = request.POST.get('Name')
    trainingbank_name_del = TrainingBank.objects.get(TrainingBankID=TrainingBankID)
    trainingbank_name_del.delete()
    return JsonResponse({"data": "true"})


# 修改教材库与角色


@csrf_exempt
def modify_trainingbank(request):
    TrainingBankID = request.POST.get('TrainingBankID')
    # inputtrainingName = request.POST.get('modifytrainingName')
    inputtrainingmaterialmaker = request.POST.get('modifytrainingmaterialmaker')
    inputtrainer = request.POST.get('modifytrainer')
    modifyTrainingBank = TrainingBank.objects.get(TrainingBankID=TrainingBankID)
    modifyTrainingBank.TrainingMaterialMaker = inputtrainingmaterialmaker
    modifyTrainingBank.Trainer = inputtrainer
    modifyTrainingBank.save()
    return JsonResponse({"data": "true"})


# 管理员界面添加教材库与角色


@csrf_exempt
def add_trainingbank(request):
    site = request.POST.get('newtrainingsite')
    plant = request.POST.get('newtrainingplant')
    input_trainingName = request.POST.get('newtrainingName')
    input_newtrainingmaterialmaker = request.POST.get('newtrainingmaterialmaker')
    input_newtrainer = request.POST.get('newtrainer')
    t = TrainingBank.objects.filter(Site=site, Plant=plant, Name=input_trainingName)
    if len(t) > 0:
        return JsonResponse({"result": "false"})
    else:
        new_TrainingBank_ID = TrainingBank(
            Site=site,
            Plant=plant,
            Name=input_trainingName,
            TrainingMaterialMaker=input_newtrainingmaterialmaker,
            Trainer=input_newtrainer,
        )
        new_TrainingBank_ID.save()
        return JsonResponse({"result": "true"})


# 管理员界面操作
# 查询所有NG类型


@csrf_exempt
def get_all_option(request):
    site = request.POST.get("admin_site")
    plant = request.POST.get("admin_plant")
    all_option = Option.objects.filter(Site=site, PlantName=plant).values()
    all_option = list(all_option)
    return JsonResponse(all_option, safe=False)


# 删除NG类型


@csrf_exempt
def del_option(request):
    OptionID = request.POST.get('OptionID')
    option_del = Option.objects.get(OptionID=OptionID)
    option_del.delete()
    return JsonResponse({"data": "true"})


# 管理员界面添加NG类型


@csrf_exempt
def add_option(request):
    option_Site = request.POST.get('newoptionsite')
    option_Plant = request.POST.get('newoptionplant')
    input_option_StaffID = request.POST.get('newoption')
    op = Option.objects.filter(
        Site=option_Site, PlantName=option_Plant, Name=input_option_StaffID
    )
    if len(op) > 0:
        return JsonResponse({"result": "false"})
    else:
        new1option = Option(
            Site=option_Site, PlantName=option_Plant, Name=input_option_StaffID
        )
        new1option.save()
        return JsonResponse({"result": "true"})
    # else:
    #     return JsonResponse({"result":"false"})


# 试卷内考生详情


@csrf_exempt
def examieed_data(request, ExamPaperID):
    # print('get_all_examieed')
    # print('ExamPaperIDYYYYY;',ExamPaperID)

    all_examieed = (
        Examinee.objects.filter(ExamPaperID=ExamPaperID).values().order_by('IsPassed')
    )
    all_examieed = list(all_examieed)
    return JsonResponse(all_examieed, safe=False)


# 删除试卷与对应考生名单


@csrf_exempt
def del_issue_paper(request):
    issue_paper_del = request.POST.get('ExamPaperID')
    issuepaper_del = ExamPaper.objects.filter(ExamPaperID=issue_paper_del)
    issuepaper_del.delete()
    examieed_del = Examinee.objects.filter(ExamPaperID=issue_paper_del)
    examieed_del.delete()
    return JsonResponse({"data": "true"})
